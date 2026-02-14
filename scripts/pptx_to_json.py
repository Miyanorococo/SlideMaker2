#!/usr/bin/env python3
"""PPTX to JSON Converter - Extract PowerPoint content to JSON format

Known Limitations:
- Connectors (LINE shapes): Extracted as simple lines with start/end points
  - Elbow connectors are converted to straight lines
  - Arrow head types (begin/end) are not preserved
  - This is a python-pptx library limitation
- Complex gradients: Non-linear gradients may not preserve exact appearance
- Text effects: Some advanced text effects may not be captured
"""
import argparse
import json
import sys
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_text_color(run, theme_colors=None, color_mapping=None, is_placeholder=False):
    """Extract text color from run, converting theme colors to RGB."""
    try:
        if run.font.color and run.font.color.type == 1:  # RGB
            rgb = run.font.color.rgb
            return f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
        elif run.font.color and run.font.color.type == 2:  # SCHEME
            theme_color = run.font.color.theme_color
            
            # Map theme color enum to name
            theme_color_names = {
                1: 'lt1', 2: 'dk1', 3: 'lt2', 4: 'dk2',
                5: 'accent1', 6: 'accent2', 7: 'accent3', 8: 'accent4',
                9: 'accent5', 10: 'accent6', 13: 'tx1', 14: 'bg1'
            }
            
            color_name = theme_color_names.get(theme_color)
            
            # Apply color mapping if available
            if color_mapping and color_name in color_mapping:
                mapped_name = color_mapping[color_name]
                if theme_colors and mapped_name in theme_colors:
                    return theme_colors[mapped_name]
            
            # Direct lookup
            if theme_colors and theme_color in theme_colors:
                return theme_colors[theme_color]
            
            # Fallback theme map
            theme_map = {
                1: "#FFFFFF", 2: "#000000", 3: "#F3F3F7", 4: "#161D26",
                5: "#41B3FF", 6: "#AD5CFF", 7: "#00E500", 8: "#FF5C85",
                9: "#FF693C", 10: "#FBD332", 13: "#000000"  # TEXT_1 default
            }
            return theme_map.get(theme_color, "#000000")
        elif run.font.color is None or run.font.color.type is None:
            # No color specified - use default
            if color_mapping:
                # Use tx1 (text color 1) from color mapping
                tx1_mapped = color_mapping.get('tx1', 'dk1')
                if theme_colors and tx1_mapped in theme_colors:
                    return theme_colors[tx1_mapped]
            return None  # Use default
    except:
        pass
    return None

def extract_theme_colors_and_mapping(pptx_path, slide_master_idx):
    """Extract actual RGB values from theme file and color mapping."""
    theme_colors = {}
    color_mapping = {}
    theme_styles = {"fill": [], "line": []}
    
    try:
        with zipfile.ZipFile(pptx_path, 'r') as zip_ref:
            # Read appropriate theme file for this master
            theme_file = f'ppt/theme/theme{slide_master_idx + 1}.xml'
            theme_xml = zip_ref.read(theme_file)
            root = ET.fromstring(theme_xml)
            
            # Define namespace
            ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main', 
                  'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
            
            # Extract color scheme
            clr_scheme = root.find('.//a:clrScheme', ns)
            if clr_scheme is not None:
                color_map = {
                    'dk1': 2,      # TEXT_1
                    'lt1': 1,      # BACKGROUND_1
                    'dk2': 4,      # TEXT_2
                    'lt2': 3,      # BACKGROUND_2
                    'accent1': 5,
                    'accent2': 6,
                    'accent3': 7,
                    'accent4': 8,
                    'accent5': 9,
                    'accent6': 10,
                }
                
                for color_name, theme_id in color_map.items():
                    color_elem = clr_scheme.find(f'.//a:{color_name}', ns)
                    if color_elem is not None:
                        srgb = color_elem.find('.//a:srgbClr', ns)
                        if srgb is not None:
                            rgb_val = srgb.get('val')
                            theme_colors[color_name] = f"#{rgb_val}"
                            theme_colors[theme_id] = f"#{rgb_val}"
                
                # Add variant for BACKGROUND_1
                if 1 in theme_colors:
                    theme_colors[14] = theme_colors[1]
            
            # Extract fill styles as raw XML strings
            fill_style_lst = root.find('.//a:fillStyleLst', ns)
            if fill_style_lst is not None:
                for fill_elem in fill_style_lst:
                    theme_styles["fill"].append(ET.tostring(fill_elem, encoding='unicode'))
            
            # Extract line styles as raw XML strings
            ln_style_lst = root.find('.//a:lnStyleLst', ns)
            if ln_style_lst is not None:
                for ln_elem in ln_style_lst:
                    theme_styles["line"].append(ET.tostring(ln_elem, encoding='unicode'))
            
            # Read slide master color mapping
            master_file = f'ppt/slideMasters/slideMaster{slide_master_idx + 1}.xml'
            master_xml = zip_ref.read(master_file)
            master_root = ET.fromstring(master_xml)
            
            clr_map = master_root.find('.//p:clrMap', ns)
            if clr_map is not None:
                # Extract mapping (e.g., bg1="dk1" means bg1 maps to dk1)
                for attr in ['bg1', 'tx1', 'bg2', 'tx2']:
                    mapped_to = clr_map.get(attr)
                    if mapped_to:
                        color_mapping[attr] = mapped_to
                        
    except Exception as e:
        print(f"Warning: Could not extract theme colors: {e}", file=sys.stderr)
    
    return theme_colors, color_mapping, theme_styles

def extract_text_with_styles(text_frame):
    """Extract text with inline styles from text frame."""
    parts = []
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            text = run.text
            if not text:
                continue
            
            # Check for hyperlink
            if run.hyperlink and run.hyperlink.address:
                parts.append(f"{{{{link:{run.hyperlink.address}:{text}}}}}")
            else:
                # Build style attributes
                styles = []
                if run.font.bold:
                    styles.append("bold")
                if run.font.italic:
                    styles.append("italic")
                if run.font.size:
                    pt = int(run.font.size.pt)
                    # Only include if not default size
                    if pt != 18:  # Skip default size
                        styles.append(f"{pt}pt")
                
                # Try to get color
                try:
                    if run.font.color and run.font.color.type == 1:  # RGB color
                        rgb = run.font.color.rgb
                        hex_color = f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
                        # Skip black color (default)
                        if hex_color != "#000000":
                            styles.append(hex_color)
                except:
                    pass
                
                if styles:
                    parts.append(f"{{{{{','.join(styles)}:{text}}}}}")
                else:
                    parts.append(text)
    
    return ''.join(parts)

def extract_line_dash(shape):
    """Extract line dash type from shape XML."""
    try:
        import xml.etree.ElementTree as ET
        sp_pr = shape._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}spPr')
        if sp_pr is not None:
            ln = sp_pr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}ln')
            if ln is not None:
                prst_dash = ln.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}prstDash')
                if prst_dash is not None:
                    dash_val = prst_dash.get('val')
                    if dash_val and dash_val != 'solid':
                        return dash_val
    except:
        pass
    return None

def extract_line_element(shape, theme_colors=None, color_mapping=None):
    """Extract line/connector as element dict."""
    try:
        elem = {
            "type": "line",
            "x": round(shape.left / 6350),
            "y": round(shape.top / 6350),
            "width": round(shape.width / 6350),
            "height": round(shape.height / 6350),
        }
        
        # Extract rotation
        if shape.rotation != 0:
            elem["rotation"] = round(shape.rotation, 1)
        
        # Extract flip from XML
        try:
            xfrm = shape._element.spPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm')
            if xfrm is not None:
                flip_h = xfrm.get('flipH')
                flip_v = xfrm.get('flipV')
                if flip_h == '1':
                    elem["flipH"] = True
                if flip_v == '1':
                    elem["flipV"] = True
        except:
            pass
        
        # Extract connector type from XML
        try:
            sp_pr = shape._element.spPr
            prst_geom = sp_pr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
            if prst_geom is not None:
                prst = prst_geom.get('prst')
                if prst:
                    # Save exact preset type
                    elem["preset"] = prst
                    
                    # Map to general connector type
                    if 'straight' in prst.lower():
                        elem["connectorType"] = "straight"
                    elif 'bent' in prst.lower():
                        elem["connectorType"] = "elbow"
                    elif 'curved' in prst.lower():
                        elem["connectorType"] = "curved"
                
                # Extract adjustments
                av_lst = prst_geom.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
                if av_lst is not None:
                    adjustments = []
                    for gd in av_lst.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}gd'):
                        fmla = gd.get('fmla', '')
                        if fmla.startswith('val '):
                            adj_val = int(fmla.split()[1])
                            adjustments.append(adj_val / 100000.0)
                    if adjustments:
                        elem["adjustments"] = adjustments
        except:
            elem["connectorType"] = "straight"  # default
        
        # Extract arrow heads from XML
        try:
            ln = shape._element.spPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}ln')
            if ln is not None:
                head_end = ln.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}headEnd')
                if head_end is not None:
                    head_type = head_end.get('type')
                    if head_type:
                        elem["headEnd"] = head_type
                
                tail_end = ln.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}tailEnd')
                if tail_end is not None:
                    tail_type = tail_end.get('type')
                    if tail_type:
                        elem["tailEnd"] = tail_type
        except:
            pass
        
        # Extract line color or gradient
        try:
            if shape.line.fill.type == 1:  # SOLID
                if shape.line.color.type == 1:  # RGB
                    rgb = shape.line.color.rgb
                    elem["color"] = f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
                elif shape.line.color.type == 2:  # SCHEME
                    theme_color = shape.line.color.theme_color
                    
                    # Map theme color to name and apply color mapping
                    theme_color_names = {
                        1: 'lt1', 2: 'dk1', 13: 'tx1'
                    }
                    color_name = theme_color_names.get(theme_color)
                    
                    if color_mapping and color_name in color_mapping:
                        mapped_name = color_mapping[color_name]
                        if theme_colors and mapped_name in theme_colors:
                            elem["color"] = theme_colors[mapped_name]
                    elif theme_colors and theme_color in theme_colors:
                        elem["color"] = theme_colors[theme_color]
                    else:
                        theme_map = {
                            1: "#FFFFFF", 2: "#000000", 13: "#000000"
                        }
                        elem["color"] = theme_map.get(theme_color, "#000000")
                
                if shape.line.width and shape.line.width > 0:
                    elem["lineWidth"] = round(shape.line.width.pt, 1)
            elif shape.line.fill.type == 3:  # GRADIENT
                # Extract line gradient
                try:
                    gradient_stops = []
                    for stop in shape.line.fill.gradient_stops:
                        stop_info = {"position": round(stop.position, 3)}
                        
                        if stop.color.type == 1:  # RGB
                            rgb = stop.color.rgb
                            stop_info["color"] = f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
                        elif stop.color.type == 2:  # SCHEME
                            theme_color = stop.color.theme_color
                            
                            # Apply color mapping
                            theme_color_names = {13: 'tx1'}
                            color_name = theme_color_names.get(theme_color)
                            
                            if color_mapping and color_name in color_mapping:
                                mapped_name = color_mapping[color_name]
                                if theme_colors and mapped_name in theme_colors:
                                    stop_info["color"] = theme_colors[mapped_name]
                            elif theme_colors and theme_color in theme_colors:
                                stop_info["color"] = theme_colors[theme_color]
                            else:
                                theme_map = {5: "#41B3FF", 6: "#AD5CFF", 7: "#00E500", 13: "#000000"}
                                stop_info["color"] = theme_map.get(theme_color, "#8B8BF8")
                        
                        gradient_stops.append(stop_info)
                    
                    if gradient_stops:
                        # Try to get angle
                        angle = 0
                        try:
                            angle = round(shape.line.fill.gradient_angle, 1)
                        except:
                            pass
                        
                        elem["lineGradient"] = {
                            "stops": gradient_stops,
                            "angle": angle
                        }
                except:
                    pass
                
                if shape.line.width and shape.line.width > 0:
                    elem["lineWidth"] = round(shape.line.width.pt, 1)
        except:
            pass
        
        return elem
    except Exception as e:
        print(f"Warning: Failed to extract line: {e}", file=sys.stderr)
        return None


def extract_freeform_element(shape, theme_colors=None, color_mapping=None):
    """Extract freeform/curve shape as element dict with path commands in px."""
    EMU_PER_PX = 6350
    try:
        sp_pr = shape._element.spPr
        cust_geom = sp_pr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}custGeom')
        if cust_geom is None:
            return None

        ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        elem = {
            "type": "freeform",
            "x": round(shape.left / EMU_PER_PX),
            "y": round(shape.top / EMU_PER_PX),
            "width": round(shape.width / EMU_PER_PX),
            "height": round(shape.height / EMU_PER_PX),
        }

        # Extract path commands
        path_el = cust_geom.find('.//a:pathLst/a:path', ns)
        if path_el is None:
            return None
        path_w = int(path_el.get('w', shape.width))
        path_h = int(path_el.get('h', shape.height))
        sx = shape.width / path_w if path_w else 1
        sy = shape.height / path_h if path_h else 1

        def to_px(x, y):
            return round(int(x) * sx / EMU_PER_PX, 1), round(int(y) * sy / EMU_PER_PX, 1)

        path = []
        for child in path_el:
            tag = child.tag.split('}')[-1]
            pts = child.findall('a:pt', ns)
            if tag == 'moveTo' and pts:
                px, py = to_px(pts[0].get('x'), pts[0].get('y'))
                path.append({"cmd": "M", "x": px, "y": py})
            elif tag == 'lnTo' and pts:
                px, py = to_px(pts[0].get('x'), pts[0].get('y'))
                path.append({"cmd": "L", "x": px, "y": py})
            elif tag == 'cubicBezTo' and len(pts) == 3:
                coords = [to_px(p.get('x'), p.get('y')) for p in pts]
                path.append({"cmd": "C", "pts": [[c[0], c[1]] for c in coords]})
            elif tag == 'close':
                path.append({"cmd": "Z"})
        elem["path"] = path

        # Fill
        try:
            no_fill = sp_pr.find('a:noFill', ns)
            if no_fill is not None:
                elem["fill"] = "none"
            elif shape.fill.type == 1:
                if shape.fill.fore_color.type == 1:
                    rgb = shape.fill.fore_color.rgb
                    elem["fill"] = f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
        except:
            pass

        # Line
        try:
            ln = sp_pr.find('a:ln', ns)
            if ln is not None:
                w = ln.get('w')
                if w:
                    elem["lineWidth"] = round(int(w) / 12700, 1)
                solid = ln.find('a:solidFill', ns)
                if solid is not None:
                    srgb = solid.find('a:srgbClr', ns)
                    scheme = solid.find('a:schemeClr', ns)
                    if srgb is not None:
                        elem["line"] = f"#{srgb.get('val')}"
                    elif scheme is not None:
                        resolved = _resolve_scheme_color(scheme.get('val'), theme_colors, color_mapping)
                        if resolved:
                            elem["line"] = resolved
        except:
            pass

        return elem
    except Exception as e:
        print(f"Warning: Failed to extract freeform: {e}", file=sys.stderr)
        return None


def extract_freeform_as_svg(shape, output_dir, slide_idx, freeform_idx):
    """Extract Freeform shape as SVG image."""
    try:
        import xml.etree.ElementTree as ET
        
        # Get custom geometry from XML
        sp_pr = shape._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}spPr')
        if sp_pr is None:
            return None
        
        cust_geom = sp_pr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}custGeom')
        if cust_geom is None:
            return None
        
        # Get shape dimensions
        width_emu = shape.width
        height_emu = shape.height
        
        # Create SVG
        svg_parts = [
            f'<svg xmlns="http://www.w3.org/2000/svg" width="{width_emu//360}" height="{height_emu//360}" viewBox="0 0 {width_emu} {height_emu}">',
            '<path d="'
        ]
        
        # Extract path commands
        path_lst = cust_geom.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}pathLst')
        if path_lst is not None:
            for path in path_lst.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}path'):
                # Get path commands (moveTo, lineTo, cubicBezTo, close, etc.)
                for child in path:
                    tag = child.tag.split('}')[-1]
                    
                    if tag == 'moveTo':
                        pt = child.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}pt')
                        if pt is not None:
                            x = pt.get('x', '0')
                            y = pt.get('y', '0')
                            svg_parts.append(f'M {x} {y} ')
                    elif tag == 'lnTo':
                        pt = child.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}pt')
                        if pt is not None:
                            x = pt.get('x', '0')
                            y = pt.get('y', '0')
                            svg_parts.append(f'L {x} {y} ')
                    elif tag == 'close':
                        svg_parts.append('Z ')
        
        svg_parts.append('" fill="none" stroke="black" stroke-width="1"/>')
        svg_parts.append('</svg>')
        
        # Save SVG
        svg_content = ''.join(svg_parts)
        svg_filename = f"slide{slide_idx + 1}_freeform{freeform_idx + 1}.svg"
        svg_path = Path(output_dir) / "images" / svg_filename
        
        svg_path.parent.mkdir(exist_ok=True)
        svg_path.write_text(svg_content)
        
        return f"images/{svg_filename}"
    except Exception as e:
        print(f"Warning: Failed to extract freeform as SVG: {e}", file=sys.stderr)
        return None

def apply_color_transforms(base_color_hex, transforms):
    """Apply color transforms (lumMod, tint) to a base color."""
    hex_color = base_color_hex.lstrip("#")
    r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
    
    has_lum_mod = 'lumMod' in transforms
    has_tint = 'tint' in transforms
    
    if has_lum_mod and has_tint:
        lum_mod_val = int(transforms['lumMod']) / 100000.0
        tint_val = int(transforms['tint']) / 100000.0
        effective_tint = tint_val / lum_mod_val
        
        r = int(round(r + ((255 - r) * effective_tint)))
        g = int(round(g + ((255 - g) * effective_tint)))
        b = int(round(b + ((255 - b) * effective_tint)))
    
    r = max(0, min(255, r))
    g = max(0, min(255, g))
    b = max(0, min(255, b))
    
    return f"#{r:02X}{g:02X}{b:02X}"

def parse_gradient_from_style(fill_style_xml, theme_colors, placeholder_color=None):
    """Parse gradient from theme fill style XML."""
    try:
        ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        fill_elem = ET.fromstring(fill_style_xml)
        
        if fill_elem.tag.endswith('gradFill'):
            grad_fill = fill_elem
        else:
            grad_fill = fill_elem.find('.//a:gradFill', ns)
        
        if grad_fill is None:
            return None
        
        gradient = {"stops": [], "angle": 0}
        
        gs_lst = grad_fill.find('.//a:gsLst', ns)
        if gs_lst is not None:
            for gs in gs_lst.findall('.//a:gs', ns):
                pos = int(gs.get('pos', 0)) / 100000.0
                
                color = "#000000"
                scheme_clr = gs.find('.//a:schemeClr', ns)
                if scheme_clr is not None:
                    clr_val = scheme_clr.get('val')
                    if clr_val == 'phClr':
                        if placeholder_color and placeholder_color in theme_colors:
                            color = theme_colors[placeholder_color]
                        else:
                            color = theme_colors.get('dk1', '#000000')
                    elif clr_val in theme_colors:
                        color = theme_colors[clr_val]
                    
                    transforms = {}
                    transform_order = ['lumMod', 'lumOff', 'satMod', 'satOff', 'tint', 'shade']
                    for tag_name in transform_order:
                        for child in scheme_clr:
                            tag = child.tag.split('}')[-1]
                            if tag == tag_name:
                                transforms[tag] = child.get('val')
                                break
                    
                    if transforms:
                        color = apply_color_transforms(color, transforms)
                
                gradient["stops"].append({"position": pos, "color": color})
        
        return gradient
    except:
        return None

def extract_shape_element(shape, theme_colors=None, color_mapping=None, theme_styles=None):
    """Extract shape as element dict."""
    try:
        # For freeform shapes, extract as rounded_rectangle (custom geometry not supported)
        shape_name = "rounded_rectangle"  # default
        
        if shape.shape_type == 5:  # MSO_SHAPE_TYPE.FREEFORM
            shape_name = "rounded_rectangle"  # Approximate as rounded rectangle
        
        # Try to get auto_shape_type
        try:
            if hasattr(shape, 'auto_shape_type'):
                from pptx.enum.shapes import MSO_SHAPE
                ast = shape.auto_shape_type
                
                # Map MSO_SHAPE to our shape names
                shape_type_map = {
                    MSO_SHAPE.RECTANGLE: "rectangle",
                    MSO_SHAPE.ROUNDED_RECTANGLE: "rounded_rectangle",
                    MSO_SHAPE.OVAL: "oval",
                    MSO_SHAPE.RIGHT_ARROW: "arrow_right",
                    MSO_SHAPE.LEFT_ARROW: "arrow_left",
                    MSO_SHAPE.UP_ARROW: "arrow_up",
                    MSO_SHAPE.DOWN_ARROW: "arrow_down",
                    MSO_SHAPE.ISOSCELES_TRIANGLE: "triangle",
                    MSO_SHAPE.DIAMOND: "diamond",
                    MSO_SHAPE.PENTAGON: "pentagon",
                    MSO_SHAPE.HEXAGON: "hexagon",
                    MSO_SHAPE.RIGHT_BRACE: "right_brace",
                    MSO_SHAPE.LEFT_BRACE: "left_brace",
                    60: "arrow_circular",  # CIRCULAR_ARROW
                }
                
                if ast in shape_type_map:
                    shape_name = shape_type_map[ast]
        except:
            # If auto_shape_type fails, check XML for preset
            try:
                    sp_pr = shape._element.spPr
                    prst_geom = sp_pr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
                    if prst_geom is not None:
                        prst = prst_geom.get('prst')
                        # Map common presets
                        preset_map = {
                            'roundRect': 'rounded_rectangle',
                            'rect': 'rectangle',
                            'ellipse': 'oval',
                            'triangle': 'triangle',
                            'diamond': 'diamond',
                            'pentagon': 'pentagon',
                            'hexagon': 'hexagon',
                        }
                        if prst in preset_map:
                            shape_name = preset_map[prst]
            except:
                pass
        
        elem = {
            "type": "shape",
            "x": round(shape.left / 6350),
            "y": round(shape.top / 6350),
            "width": round(shape.width / 6350),
            "height": round(shape.height / 6350),
            "shape": shape_name
        }
        
        # Check for style references
        style_fill_idx = None
        style_fill_color = None
        
        try:
            sp_elem = shape._element
            style = sp_elem.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}style')
            if style is not None:
                fill_ref = style.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}fillRef')
                if fill_ref is not None:
                    style_fill_idx = int(fill_ref.get('idx', 0))
                    scheme_clr = fill_ref.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}schemeClr')
                    if scheme_clr is not None:
                        style_fill_color = scheme_clr.get('val')
        except:
            pass
        
        # Extract rotation
        if shape.rotation != 0:
            elem["rotation"] = round(shape.rotation, 1)
        
        # Extract adjustments (shape control points)
        if hasattr(shape, 'adjustments') and shape.adjustments:
            try:
                adjustments = []
                for i in range(len(shape.adjustments)):
                    adj_val = shape.adjustments[i]
                    if adj_val is not None:
                        adjustments.append(round(adj_val, 5))
                if adjustments:
                    elem["adjustments"] = adjustments
            except:
                # If adjustments fail, try to get from XML
                try:
                    sp_pr = shape._element.spPr
                    prst_geom = sp_pr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
                    if prst_geom is not None:
                        av_lst = prst_geom.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
                        if av_lst is not None and len(av_lst) > 0:
                            # Has custom adjustments in XML
                            pass  # Use default for now
                        else:
                            # No adjustments specified, use default
                            # For roundRect, default is typically 0.16667
                            if shape_name == "rounded_rectangle":
                                elem["adjustments"] = [0.16667]
                except:
                    pass
        
        # Extract fill color
        try:
            has_explicit_fill = shape.fill.type is not None and shape.fill.type != 0
            
            # Resolve style reference if no explicit fill
            if not has_explicit_fill and style_fill_idx is not None and theme_styles and theme_colors:
                if 0 < style_fill_idx <= len(theme_styles["fill"]):
                    fill_style_xml = theme_styles["fill"][style_fill_idx - 1]
                    gradient = parse_gradient_from_style(fill_style_xml, theme_colors, style_fill_color)
                    if gradient:
                        elem["gradient"] = gradient
            
            if shape.fill.type == 1:  # SOLID
                if shape.fill.fore_color.type == 1:  # RGB
                    rgb = shape.fill.fore_color.rgb
                    elem["fill"] = f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
                elif shape.fill.fore_color.type == 2:  # SCHEME (theme color)
                    theme_color = shape.fill.fore_color.theme_color
                    
                    # Get the scheme color name (e.g., "bg1", "accent1")
                    # theme_color is an enum value, need to map to name
                    theme_color_names = {
                        1: 'lt1',      # BACKGROUND_1
                        2: 'dk1',      # TEXT_1
                        3: 'lt2',      # BACKGROUND_2
                        4: 'dk2',      # TEXT_2
                        5: 'accent1',
                        6: 'accent2',
                        7: 'accent3',
                        8: 'accent4',
                        9: 'accent5',
                        10: 'accent6',
                        14: 'bg1',     # BACKGROUND_1 variant (used in shapes)
                    }
                    
                    color_name = theme_color_names.get(theme_color, f'color{theme_color}')
                    
                    # Apply color mapping if available (e.g., bg1 -> dk1)
                    if color_mapping and color_name in color_mapping:
                        mapped_name = color_mapping[color_name]
                        # Get the actual color from theme
                        if theme_colors and mapped_name in theme_colors:
                            elem["fill"] = theme_colors[mapped_name]
                        elif theme_colors and theme_color in theme_colors:
                            elem["fill"] = theme_colors[theme_color]
                    elif theme_colors and theme_color in theme_colors:
                        elem["fill"] = theme_colors[theme_color]
                    else:
                        # Fallback to default mapping
                        theme_map = {
                            1: "#FFFFFF",  # BACKGROUND_1 (white)
                            2: "#000000",  # TEXT_1 (black)
                            3: "#F3F3F7",  # BACKGROUND_2 (light gray)
                            4: "#161D26",  # TEXT_2 (dark gray)
                            5: "#41B3FF",  # ACCENT_1 (light blue)
                            6: "#AD5CFF",  # ACCENT_2 (purple)
                            7: "#00E500",  # ACCENT_3 (green)
                            8: "#FF5C85",  # ACCENT_4 (pink)
                            9: "#FF693C",  # ACCENT_5 (orange)
                            10: "#FBD332", # ACCENT_6 (yellow)
                            14: "#FFFFFF", # BACKGROUND_1 (white)
                        }
                        if theme_color in theme_map:
                            elem["fill"] = theme_map[theme_color]
                
                # Extract transparency/alpha
                try:
                    import xml.etree.ElementTree as ET
                    sp_pr = shape._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}spPr')
                    if sp_pr is not None:
                        solid_fill = sp_pr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
                        if solid_fill is not None:
                            alpha = solid_fill.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
                            if alpha is not None:
                                alpha_val = int(alpha.get('val', 100000))
                                # Convert to percentage (100000 = 100% opaque, 0 = 100% transparent)
                                opacity = round(alpha_val / 1000, 1)  # 0-100
                                if opacity < 100:
                                    elem["opacity"] = opacity
                except:
                    pass
            elif shape.fill.type == 3:  # GRADIENT
                # Extract gradient information
                try:
                    gradient_stops = []
                    
                    # Get gradient stops with alpha from XML
                    import xml.etree.ElementTree as ET
                    sp_pr = shape._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}spPr')
                    grad_fill_xml = None
                    if sp_pr is not None:
                        grad_fill_xml = sp_pr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}gradFill')
                    
                    for i, stop in enumerate(shape.fill.gradient_stops):
                        stop_info = {"position": round(stop.position, 3)}
                        
                        if stop.color.type == 1:  # RGB
                            rgb = stop.color.rgb
                            stop_info["color"] = f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
                        elif stop.color.type == 2:  # SCHEME
                            theme_color = stop.color.theme_color
                            if theme_colors and theme_color in theme_colors:
                                stop_info["color"] = theme_colors[theme_color]
                            else:
                                # Fallback
                                theme_map = {5: "#41B3FF", 6: "#AD5CFF", 7: "#00E500"}
                                stop_info["color"] = theme_map.get(theme_color, "#41B3FF")
                        
                        # Extract alpha from XML
                        if grad_fill_xml is not None:
                            try:
                                gs_lst = grad_fill_xml.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}gsLst')
                                if gs_lst is not None:
                                    gs_elements = gs_lst.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}gs')
                                    if i < len(gs_elements):
                                        alpha = gs_elements[i].find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
                                        if alpha is not None:
                                            alpha_val = int(alpha.get('val', 100000))
                                            opacity = round(alpha_val / 1000, 1)
                                            if opacity < 100:
                                                stop_info["opacity"] = opacity
                            except:
                                pass
                        
                        gradient_stops.append(stop_info)
                    
                    if gradient_stops:
                        # Try to get angle (only for linear gradients)
                        angle = 0
                        try:
                            angle = round(shape.fill.gradient_angle, 1)
                        except:
                            # Non-linear gradient (radial, path, etc.) - use 0
                            pass
                        
                        elem["gradient"] = {
                            "stops": gradient_stops,
                            "angle": angle
                        }
                except Exception as e:
                    # Fallback: use representative color
                    elem["fill"] = "#41B3FF"
            elif shape.fill.type is None or shape.fill.type == 0:  # NO_FILL
                elem["fill"] = "none"
        except:
            pass
        
        # Extract line color
        try:
            if shape.line.fill.type == 1:  # SOLID
                if shape.line.color.type == 1:  # RGB
                    rgb = shape.line.color.rgb
                    elem["line"] = f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
                elif shape.line.color.type == 2:  # SCHEME
                    theme_color = shape.line.color.theme_color
                    
                    # Apply color mapping
                    color_name = {1: 'lt1', 2: 'dk1', 3: 'lt2', 4: 'dk2', 5: 'accent1', 6: 'accent2', 7: 'accent3', 8: 'accent4', 9: 'accent5', 10: 'accent6', 13: 'tx1', 14: 'bg1', 15: 'tx2'}.get(theme_color)
                    
                    if color_mapping and color_name and color_name in color_mapping:
                        mapped_name = color_mapping[color_name]
                        if theme_colors and mapped_name in theme_colors:
                            elem["line"] = theme_colors[mapped_name]
                        elif theme_colors and theme_color in theme_colors:
                            elem["line"] = theme_colors[theme_color]
                    elif theme_colors and theme_color in theme_colors:
                        elem["line"] = theme_colors[theme_color]
                    else:
                        elem["line"] = {1: "#FFFFFF", 2: "#000000", 13: "#000000", 15: "#161D26"}.get(theme_color, "#000000")
                
                if shape.line.width and shape.line.width > 0:
                    elem["lineWidth"] = round(shape.line.width.pt, 1)
                
                # Extract dash type
                dash = extract_line_dash(shape)
                if dash:
                    elem["lineDash"] = dash
            elif shape.line.fill.type == 3:  # GRADIENT
                # Extract line gradient
                try:
                    gradient_stops = []
                    for stop in shape.line.fill.gradient_stops:
                        stop_info = {"position": round(stop.position, 3)}
                        
                        if stop.color.type == 1:  # RGB
                            rgb = stop.color.rgb
                            stop_info["color"] = f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
                        elif stop.color.type == 2:  # SCHEME
                            theme_color = stop.color.theme_color
                            if theme_colors and theme_color in theme_colors:
                                stop_info["color"] = theme_colors[theme_color]
                            else:
                                theme_map = {5: "#41B3FF", 6: "#AD5CFF", 7: "#00E500"}
                                stop_info["color"] = theme_map.get(theme_color, "#8B8BF8")
                        
                        gradient_stops.append(stop_info)
                    
                    if gradient_stops:
                        # Try to get angle, but don't fail if not available
                        angle = 0
                        try:
                            angle = round(shape.line.fill.gradient_angle, 1)
                        except:
                            pass
                        
                        elem["lineGradient"] = {
                            "stops": gradient_stops,
                            "angle": angle
                        }
                except Exception as e:
                    # Fallback
                    elem["line"] = "#8B8BF8"
                
                if shape.line.width and shape.line.width > 0:
                    elem["lineWidth"] = round(shape.line.width.pt, 1)
                
                # Extract dash type
                dash = extract_line_dash(shape)
                if dash:
                    elem["lineDash"] = dash
            elif shape.line.fill.type == 5 or shape.line.fill.type is None or shape.line.fill.type == 0:  # BACKGROUND or NO_FILL
                # Check if there's explicit noFill in XML
                has_explicit_no_fill = False
                try:
                    import xml.etree.ElementTree as ET
                    sp_pr = shape._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}spPr')
                    if sp_pr is not None:
                        ln = sp_pr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}ln')
                        if ln is not None:
                            no_fill = ln.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}noFill')
                            if no_fill is not None:
                                has_explicit_no_fill = True
                except:
                    pass
                
                if has_explicit_no_fill:
                    # Explicit noFill - no line
                    elem["line"] = "none"
                else:
                    # Check if there's a style reference in XML
                    try:
                        import xml.etree.ElementTree as ET
                        style = shape._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}style')
                        if style is not None:
                            ln_ref = style.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}lnRef')
                            if ln_ref is not None:
                                idx = ln_ref.get('idx')
                                scheme_clr = ln_ref.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}schemeClr')
                                
                                if scheme_clr is not None and idx and int(idx) > 0:
                                    # Has line style reference
                                    scheme_val = scheme_clr.get('val')
                                    
                                    # Map scheme color name to theme color
                                    scheme_map = {
                                        'accent1': 5, 'accent2': 6, 'accent3': 7,
                                        'accent4': 8, 'accent5': 9, 'accent6': 10,
                                        'dk1': 2, 'lt1': 1, 'dk2': 4, 'lt2': 3
                                    }
                                    
                                    theme_color_id = scheme_map.get(scheme_val)
                                    if theme_color_id and theme_colors and theme_color_id in theme_colors:
                                        elem["line"] = theme_colors[theme_color_id]
                                        
                                        # Get line width from theme (idx 1=6350, 2=12700, 3=19050 EMU)
                                        width_map = {1: 0.5, 2: 1.0, 3: 1.5}
                                        elem["lineWidth"] = width_map.get(int(idx), 1.0)
                                        
                                        # Extract dash type
                                        dash = extract_line_dash(shape)
                                        if dash:
                                            elem["lineDash"] = dash
                                    else:
                                        elem["line"] = "none"
                                else:
                                    elem["line"] = "none"
                            else:
                                elem["line"] = "none"
                        else:
                            elem["line"] = "none"
                    except:
                        elem["line"] = "none"
        except:
            pass
        
        # Extract text with styles
        if shape.has_text_frame and shape.text.strip():
            # Extract text frame margins
            tf = shape.text_frame
            if tf.margin_left and tf.margin_left != 91440:
                elem["marginLeft"] = tf.margin_left
            if tf.margin_top and tf.margin_top != 45720:
                elem["marginTop"] = tf.margin_top
            if tf.margin_right and tf.margin_right != 91440:
                elem["marginRight"] = tf.margin_right
            if tf.margin_bottom and tf.margin_bottom != 45720:
                elem["marginBottom"] = tf.margin_bottom
            
            # Extract vertical anchor
            if tf.vertical_anchor is not None:
                elem["verticalAnchor"] = int(tf.vertical_anchor)
            
            text_parts = []
            default_font_size = None
            text_align = None
            
            # Get text alignment from first paragraph
            if shape.text_frame.paragraphs:
                para = shape.text_frame.paragraphs[0]
                if para.alignment is not None:
                    if para.alignment == 2:  # CENTER
                        text_align = "center"
                    elif para.alignment == 3:  # RIGHT
                        text_align = "right"
                    elif para.alignment == 1:  # LEFT
                        text_align = "left"
            
            # Check if multiple paragraphs (bullet list)
            paragraphs_with_text = [p for p in shape.text_frame.paragraphs if p.text.strip()]
            
            # Check if any paragraph has bullet/numbering
            has_bullets = False
            for para in paragraphs_with_text:
                try:
                    pPr = para._element.pPr
                    if pPr is not None:
                        ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                        bu_char = pPr.find('.//a:buChar', ns)
                        bu_auto_num = pPr.find('.//a:buAutoNum', ns)
                        if bu_char is not None or bu_auto_num is not None:
                            has_bullets = True
                            break
                except:
                    pass
            
            if has_bullets:
                # Has bullets - extract as items array (even if single paragraph)
                items = []
                for paragraph in paragraphs_with_text:
                    # Get font size
                    if not default_font_size:
                        for run in paragraph.runs:
                            if run.font.size:
                                default_font_size = int(run.font.size.pt)
                                break
                    
                    parts = []
                    for run in paragraph.runs:
                        if not run.text:
                            continue
                        
                        # Check for hyperlink
                        if run.hyperlink and run.hyperlink.address:
                            parts.append(f"{{{{link:{run.hyperlink.address}:{run.text}}}}}")
                        else:
                            styles = []
                            if run.font.bold:
                                styles.append("bold")
                            if run.font.italic:
                                styles.append("italic")
                            if run.font.size:
                                pt = int(run.font.size.pt)
                                if default_font_size and pt != default_font_size:
                                    styles.append(f"{pt}pt")
                            
                            try:
                                hex_color = extract_text_color(run, theme_colors)
                                if hex_color and hex_color != "#000000":
                                    styles.append(hex_color)
                            except:
                                pass
                            
                            if styles:
                                parts.append(f"{{{{{','.join(styles)}:{run.text}}}}}")
                            else:
                                parts.append(run.text)
                    
                    item_text = ''.join(parts)
                    if item_text.strip():
                        items.append(item_text)
                
                if items:
                    elem["items"] = items
                    # Note: textAlign will be set below
            else:
                # Single paragraph or multiple paragraphs without bullets - extract as text
                text_parts = []
                for para_idx, paragraph in enumerate(paragraphs_with_text):
                    # Get font size
                    if not default_font_size:
                        for run in paragraph.runs:
                            if run.font.size:
                                default_font_size = int(run.font.size.pt)
                                break
                    
                    # Add newline between paragraphs
                    if para_idx > 0:
                        text_parts.append('\n')
                    
                    for run in paragraph.runs:
                        if not run.text:
                            continue
                        
                        # Check for hyperlink
                        if run.hyperlink and run.hyperlink.address:
                            text_parts.append(f"{{{{link:{run.hyperlink.address}:{run.text}}}}}")
                        else:
                            styles = []
                            if run.font.bold:
                                styles.append("bold")
                            if run.font.italic:
                                styles.append("italic")
                            if run.font.size:
                                pt = int(run.font.size.pt)
                                if default_font_size and pt != default_font_size:
                                    styles.append(f"{pt}pt")
                            
                            try:
                                hex_color = extract_text_color(run, theme_colors)
                                if hex_color and hex_color != "#000000":
                                    styles.append(hex_color)
                            except:
                                pass
                            
                            if styles:
                                text_parts.append(f"{{{{{','.join(styles)}:{run.text}}}}}")
                            else:
                                text_parts.append(run.text)
                
                elem["text"] = ''.join(text_parts)
            
            # Add fontSize if not default
            if default_font_size and default_font_size != 18:
                elem["fontSize"] = default_font_size
            
            # Add text alignment (always set for shapes to ensure correct alignment)
            if text_align:
                elem["textAlign"] = text_align
            else:
                # Default to left if not specified
                elem["textAlign"] = "left"
        
        # Extract hyperlink
        try:
            if hasattr(shape, 'click_action') and shape.click_action.hyperlink and shape.click_action.hyperlink.address:
                elem["link"] = shape.click_action.hyperlink.address
            else:
                # Remove null link
                if "link" in elem and elem["link"] is None:
                    del elem["link"]
        except:
            pass
        
        return elem
    except Exception as e:
        print(f"Warning: Failed to extract shape details: {e}", file=sys.stderr)
        return None

def extract_textbox_element(shape, theme_colors=None, color_mapping=None, theme_styles=None):
    """Extract textbox as element dict."""
    # Check if it's actually a shape with preset geometry (not a plain textbox)
    try:
        sp_pr = shape._element.spPr
        prst_geom = sp_pr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
        if prst_geom is not None:
            prst = prst_geom.get('prst')
            # If it has any preset geometry (not just 'rect'), treat as shape
            if prst and prst != 'rect':
                # This is a shape with text, not a plain textbox
                return extract_shape_element(shape, theme_colors, color_mapping, theme_styles)
    except:
        pass
    
    elem = {
        "type": "textbox",
        "x": round(shape.left / 6350),  # px (1920x1080 basis)
        "y": round(shape.top / 6350),
        "width": round(shape.width / 6350),
    }
    
    # Extract rotation
    if shape.rotation != 0:
        elem["rotation"] = round(shape.rotation, 1)
    
    # Extract flip
    try:
        xfrm = shape._element.spPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm')
        if xfrm is not None:
            flip_h = xfrm.get('flipH')
            flip_v = xfrm.get('flipV')
            if flip_h == '1':
                elem["flipH"] = True
            if flip_v == '1':
                elem["flipV"] = True
    except:
        pass
    
    # Extract fill color
    try:
        if shape.fill.type == 1:  # SOLID
            if shape.fill.fore_color.type == 1:  # RGB
                rgb = shape.fill.fore_color.rgb
                elem["fill"] = f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
            elif shape.fill.fore_color.type == 2:  # SCHEME
                theme_color = shape.fill.fore_color.theme_color
                if theme_colors and theme_color in theme_colors:
                    elem["fill"] = theme_colors[theme_color]
        elif shape.fill.type is None or shape.fill.type == 0:
            elem["fill"] = "none"
    except:
        pass
    
    # Extract line color
    try:
        if shape.line.fill.type == 1:  # SOLID
            if shape.line.color.type == 1:  # RGB
                rgb = shape.line.color.rgb
                elem["line"] = f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
            elif shape.line.color.type == 2:  # SCHEME
                theme_color = shape.line.color.theme_color
                if theme_colors and theme_color in theme_colors:
                    elem["line"] = theme_colors[theme_color]
                else:
                    theme_map = {
                        1: "#FFFFFF", 2: "#000000", 5: "#41B3FF", 6: "#AD5CFF",
                        7: "#00E500", 8: "#FF5C85", 9: "#FF693C", 10: "#FBD332",
                        13: "#000000"
                    }
                    elem["line"] = theme_map.get(theme_color, "#000000")
            
            if shape.line.width and shape.line.width > 0:
                elem["lineWidth"] = round(shape.line.width.pt, 1)
        elif shape.line.fill.type == 5 or shape.line.fill.type is None or shape.line.fill.type == 0:
            elem["line"] = "none"
    except:
        pass
    
def extract_textbox_element(shape, theme_colors=None, color_mapping=None, theme_styles=None):
    """Extract textbox as element dict."""
    # Check if it's actually a shape with preset geometry (not a plain textbox)
    try:
        sp_pr = shape._element.spPr
        prst_geom = sp_pr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
        if prst_geom is not None:
            prst = prst_geom.get('prst')
            # If it has any preset geometry (not just 'rect'), treat as shape
            if prst and prst != 'rect':
                # This is a shape with text, not a plain textbox
                return extract_shape_element(shape, theme_colors, color_mapping, theme_styles)
    except:
        pass
    
    elem = {
        "type": "textbox",
        "x": round(shape.left / 6350),  # px (1920x1080 basis)
        "y": round(shape.top / 6350),
        "width": round(shape.width / 6350),
    }
    
    # Extract height (for TEXT_TO_FIT_SHAPE auto-shrink)
    if shape.height:
        h_px = round(shape.height / 6350)
        if h_px > 10:
            elem["height"] = h_px
    # Extract rotation
    if shape.rotation != 0:
        elem["rotation"] = round(shape.rotation, 1)
    
    # Extract flip
    try:
        xfrm = shape._element.spPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm')
        if xfrm is not None:
            flip_h = xfrm.get('flipH')
            flip_v = xfrm.get('flipV')
            if flip_h == '1':
                elem["flipH"] = True
            if flip_v == '1':
                elem["flipV"] = True
    except:
        pass
    
    # Extract fill color
    try:
        if shape.fill.type == 1:  # SOLID
            if shape.fill.fore_color.type == 1:  # RGB
                rgb = shape.fill.fore_color.rgb
                elem["fill"] = f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
            elif shape.fill.fore_color.type == 2:  # SCHEME
                theme_color = shape.fill.fore_color.theme_color
                if theme_colors and theme_color in theme_colors:
                    elem["fill"] = theme_colors[theme_color]
        elif shape.fill.type is None or shape.fill.type == 0:
            elem["fill"] = "none"
    except:
        pass
    
    # Extract line color
    try:
        if shape.line.fill.type == 1:  # SOLID
            if shape.line.color.type == 1:  # RGB
                rgb = shape.line.color.rgb
                elem["line"] = f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
            elif shape.line.color.type == 2:  # SCHEME
                theme_color = shape.line.color.theme_color
                if theme_colors and theme_color in theme_colors:
                    elem["line"] = theme_colors[theme_color]
                else:
                    theme_map = {
                        1: "#FFFFFF", 2: "#000000", 5: "#41B3FF", 6: "#AD5CFF",
                        7: "#00E500", 8: "#FF5C85", 9: "#FF693C", 10: "#FBD332",
                        13: "#000000"
                    }
                    elem["line"] = theme_map.get(theme_color, "#000000")
            
            if shape.line.width and shape.line.width > 0:
                elem["lineWidth"] = round(shape.line.width.pt, 1)
        elif shape.line.fill.type == 5 or shape.line.fill.type is None or shape.line.fill.type == 0:
            elem["line"] = "none"
    except:
        pass
    
    # Extract text with styles
    text_parts = []
    default_font_size = None
    
    # Determine default text color from master
    default_text_color = None
    if color_mapping:
        tx1_mapped = color_mapping.get('tx1', 'dk1')
        if theme_colors and tx1_mapped in theme_colors:
            default_text_color = theme_colors[tx1_mapped]
    
    # Check if multiple paragraphs (should be items array)
    paragraphs_with_text = [p for p in shape.text_frame.paragraphs if p.text.strip()]
    all_paragraphs = list(shape.text_frame.paragraphs)
    
    if len(all_paragraphs) > 1:
        # Multiple paragraphs - extract as paragraphs with bullet info
        paragraphs = []
        for paragraph in all_paragraphs:
            # Get font size
            if not default_font_size:
                for run in paragraph.runs:
                    if run.font.size:
                        default_font_size = int(run.font.size.pt)
                        break
            
            # Empty paragraph
            if not paragraph.text.strip():
                paragraphs.append({"text": "", "bullet": False})
                continue
            
            # Check for bullet or numbering
            has_bullet = False
            numbering_type = None
            bu_font = None
            mar_l = None
            indent = None
            space_after = None
            try:
                pPr = paragraph._element.pPr
                if pPr is not None:
                    bu_auto_num = pPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}buAutoNum')
                    bu_char = pPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}buChar')
                    
                    if bu_auto_num is not None:
                        numbering_type = bu_auto_num.get('type', 'arabicPeriod')
                    elif bu_char is not None:
                        has_bullet = True
                    
                    bu_font_elem = pPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}buFont')
                    if bu_font_elem is not None:
                        bu_font = bu_font_elem.get('typeface')
                    mar_l = pPr.get('marL')
                    indent = pPr.get('indent')
                    ns_a = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                    spc_aft = pPr.find('.//a:spcAft/a:spcPts', ns_a)
                    if spc_aft is not None:
                        space_after = spc_aft.get('val')
            except:
                pass
            
            parts = []
            for run in paragraph.runs:
                if not run.text:
                    continue
                
                # Check for hyperlink
                if run.hyperlink and run.hyperlink.address:
                    parts.append(f"{{{{link:{run.hyperlink.address}:{run.text}}}}}")
                else:
                    styles = []
                    if run.font.bold:
                        styles.append("bold")
                    if run.font.italic:
                        styles.append("italic")
                    if run.font.size:
                        pt = int(run.font.size.pt)
                        if default_font_size and pt != default_font_size:
                            styles.append(f"{pt}pt")
                    
                    try:
                        hex_color = extract_text_color(run, theme_colors, color_mapping, is_placeholder=False)
                        if hex_color:
                            if hex_color != "#000000" or (default_text_color and default_text_color != "#000000"):
                                styles.append(hex_color)
                        elif default_text_color and default_text_color != "#000000":
                            styles.append(default_text_color)
                    except:
                        pass
                    
                    if styles:
                        parts.append(f"{{{{{','.join(styles)}:{run.text}}}}}")
                    else:
                        parts.append(run.text)
            
            item_text = ''.join(parts)
            para_info = {"text": item_text}
            if has_bullet:
                para_info["bullet"] = True
            elif numbering_type:
                para_info["numbering"] = numbering_type
            if bu_font:
                para_info["buFont"] = bu_font
            if mar_l:
                para_info["marL"] = int(mar_l)
            if indent:
                para_info["indent"] = int(indent)
            if space_after and space_after != "0":
                para_info["spaceAfter"] = int(space_after)
            
            paragraphs.append(para_info)
        
        if paragraphs:
            elem["paragraphs"] = paragraphs  # Use "paragraphs" instead of "items" for textbox
            
            # Add fontSize if not default
            if default_font_size and default_font_size != 18:
                elem["fontSize"] = default_font_size
            
            # Get alignment
            if paragraphs_with_text[0].alignment == 2:
                elem["align"] = "center"
            elif paragraphs_with_text[0].alignment == 3:
                elem["align"] = "right"
            
            return elem
    
    # Single paragraph - extract as text
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if not run.text:
                continue
            
            # Get font size for this run
            if run.font.size and not default_font_size:
                default_font_size = int(run.font.size.pt)
            
            # Check for hyperlink
            if run.hyperlink and run.hyperlink.address:
                text_parts.append(f"{{{{link:{run.hyperlink.address}:{run.text}}}}}")
            else:
                # Build style attributes
                styles = []
                if run.font.bold:
                    styles.append("bold")
                if run.font.italic:
                    styles.append("italic")
                
                # Try to get color (including default color from master)
                try:
                    hex_color = extract_text_color(run, theme_colors, color_mapping, is_placeholder=False)
                    if hex_color:
                        # Always include color if it's not black or if default is not black
                        if hex_color != "#000000" or (default_text_color and default_text_color != "#000000"):
                            styles.append(hex_color)
                    elif default_text_color and default_text_color != "#000000":
                        # No color specified, use default from master
                        styles.append(default_text_color)
                except:
                    pass
                
                if styles:
                    text_parts.append(f"{{{{{','.join(styles)}:{run.text}}}}}")
                else:
                    text_parts.append(run.text)
    
    elem["text"] = ''.join(text_parts)
    
    # Add fontSize if consistent and not default
    if default_font_size and default_font_size != 18:
        elem["fontSize"] = default_font_size
    
    # Detect alignment
    if shape.text_frame.paragraphs:
        align = shape.text_frame.paragraphs[0].alignment
        if align == 2:  # CENTER
            elem["align"] = "center"
        elif align == 3:  # RIGHT
            elem["align"] = "right"
        # Don't add align for left (1) or None (default)
    
    return elem

def extract_picture_element(shape, output_dir=None, slide_idx=0, img_idx=0):
    """Extract picture as element dict and save image file."""
    elem = {
        "type": "image",
        "x": round(shape.left / 6350),
        "y": round(shape.top / 6350),
        "width": round(shape.width / 6350),
        "height": round(shape.height / 6350),
    }
    
    # Save image to file
    if output_dir:
        try:
            image = shape.image
            image_bytes = image.blob
            
            # Determine format
            ext = shape.image.ext or "png"
            
            # Create images directory
            images_dir = Path(output_dir) / "images"
            images_dir.mkdir(exist_ok=True)
            
            # Save image
            image_filename = f"slide{slide_idx + 1}_image{img_idx + 1}.{ext}"
            image_path = images_dir / image_filename
            
            with open(image_path, 'wb') as f:
                f.write(image_bytes)
            
            # Store relative path
            elem["src"] = f"images/{image_filename}"
        except Exception as e:
            print(f"Warning: Failed to save image: {e}", file=sys.stderr)
    
    # Extract hyperlink
    if hasattr(shape, 'click_action') and shape.click_action.hyperlink:
        elem["link"] = shape.click_action.hyperlink.address
    
    return elem

def extract_group_element(shape, theme_colors=None, color_mapping=None, theme_styles=None, output_dir=None, slide_idx=0, img_counter=0):
    """Extract group as element dict with nested elements.
    
    Note: python-pptx returns absolute slide coordinates for grouped shapes.
    """
    elem = {
        "type": "group",
        "x": round(shape.left / 6350),
        "y": round(shape.top / 6350),
        "width": round(shape.width / 6350),
        "height": round(shape.height / 6350),
        "elements": []
    }
    
    # Extract rotation
    if shape.rotation != 0:
        elem["rotation"] = round(shape.rotation, 1)
    
    # Extract each shape in the group
    for sub_shape in shape.shapes:
        try:
            sub_elem = None
            
            # Handle nested groups recursively
            if sub_shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                sub_elem, img_counter = extract_group_element(sub_shape, theme_colors, color_mapping, theme_styles, output_dir, slide_idx, img_counter)
                # For nested groups, positions are already calculated
                if sub_elem:
                    elem["elements"].append(sub_elem)
                continue  # Skip position recalculation
            elif sub_shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                sub_elem = extract_textbox_element(sub_shape, theme_colors, color_mapping)
            elif sub_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                sub_elem = extract_picture_element(sub_shape, output_dir, slide_idx, img_counter)
                if sub_elem:
                    img_counter += 1
            elif sub_shape.shape_type in (MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM):
                sub_elem = extract_shape_element(sub_shape, theme_colors, color_mapping)
            elif sub_shape.shape_type == MSO_SHAPE_TYPE.LINE:
                sub_elem = extract_line_element(sub_shape, theme_colors, color_mapping)
            
            if sub_elem:
                # Transform coordinates from child coordinate system to slide coordinates
                import xml.etree.ElementTree as ET
                grp_sp_pr = shape._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}grpSpPr')
                xfrm = grp_sp_pr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm') if grp_sp_pr else None
                
                if xfrm is not None:
                    off = xfrm.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}off')
                    ext = xfrm.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}ext')
                    ch_off = xfrm.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}chOff')
                    ch_ext = xfrm.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}chExt')
                    
                    if off is not None and ch_off is not None and ext is not None and ch_ext is not None:
                        group_off_x = int(off.get('x'))
                        group_off_y = int(off.get('y'))
                        group_ext_cx = int(ext.get('cx'))
                        group_ext_cy = int(ext.get('cy'))
                        ch_off_x = int(ch_off.get('x'))
                        ch_off_y = int(ch_off.get('y'))
                        ch_ext_cx = int(ch_ext.get('cx'))
                        ch_ext_cy = int(ch_ext.get('cy'))
                        
                        # Transform: abs = group_off + (child - chOff) * (group_ext / ch_ext)
                        child_x = sub_shape.left
                        child_y = sub_shape.top
                        
                        scale_x = group_ext_cx / ch_ext_cx if ch_ext_cx != 0 else 1
                        scale_y = group_ext_cy / ch_ext_cy if ch_ext_cy != 0 else 1
                        
                        abs_x = group_off_x + (child_x - ch_off_x) * scale_x
                        abs_y = group_off_y + (child_y - ch_off_y) * scale_y
                        
                        sub_elem["x"] = round(abs_x / 6350)
                        sub_elem["y"] = round(abs_y / 6350)
                    else:
                        # Fallback: use python-pptx coordinates as-is
                        pass
                
                elem["elements"].append(sub_elem)
        except Exception as e:
            print(f"Warning: Failed to extract grouped shape: {e}", file=sys.stderr)
    
    return elem, img_counter

def _resolve_scheme_color(val, theme_colors, color_mapping):
    """Resolve scheme color name (from XML schemeClr val) to RGB hex."""
    # color_mapping first (e.g. tx1->dk1, bg2->lt2)
    if color_mapping and val in color_mapping:
        mapped = color_mapping[val]
        if theme_colors and mapped in theme_colors:
            return theme_colors[mapped]
    # Direct lookup by name
    if theme_colors and val in theme_colors:
        return theme_colors[val]
    # By enum id (same map as extract_text_color)
    name_to_id = {'lt1': 1, 'dk1': 2, 'lt2': 3, 'dk2': 4,
                  'accent1': 5, 'accent2': 6, 'accent3': 7, 'accent4': 8,
                  'accent5': 9, 'accent6': 10, 'tx1': 13, 'bg1': 14, 'tx2': 4, 'bg2': 3}
    tid = name_to_id.get(val)
    if tid and theme_colors and tid in theme_colors:
        return theme_colors[tid]
    # Fallback (same as extract_text_color)
    fallback = {'lt1': "#FFFFFF", 'dk1': "#000000", 'lt2': "#F3F3F7", 'dk2': "#161D26",
                'accent1': "#41B3FF", 'accent2': "#AD5CFF", 'accent3': "#00E500",
                'accent4': "#FF5C85", 'accent5': "#FF693C", 'accent6': "#FBD332",
                'tx1': "#000000", 'bg1': "#FFFFFF", 'tx2': "#161D26", 'bg2': "#F3F3F7"}
    return fallback.get(val)


def _extract_cell(cell, theme_colors=None, color_mapping=None):
    """Extract cell as string (text only) or dict (has extra properties)."""
    ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    tc = cell._tc
    tc_pr = tc.find('a:tcPr', ns)
    text = cell.text
    props = {}

    if tc_pr is not None:
        # Fill
        solid = tc_pr.find('a:solidFill', ns)
        if solid is not None:
            srgb = solid.find('a:srgbClr', ns)
            scheme = solid.find('a:schemeClr', ns)
            if srgb is not None:
                props["fill"] = f"#{srgb.get('val')}"
            elif scheme is not None:
                resolved = _resolve_scheme_color(scheme.get('val'), theme_colors, color_mapping)
                if resolved:
                    props["fill"] = resolved

        # Borders
        borders = {}
        for side, tag in [("left", "lnL"), ("right", "lnR"), ("top", "lnT"), ("bottom", "lnB")]:
            ln = tc_pr.find(f'a:{tag}', ns)
            if ln is None:
                continue
            border = {}
            w = ln.get('w')
            if w:
                border["width"] = round(int(w) / 12700, 1)  # EMU to pt
            if ln.find('a:noFill', ns) is not None:
                border["fill"] = "none"
            else:
                sf = ln.find('a:solidFill', ns)
                if sf is not None:
                    srgb = sf.find('a:srgbClr', ns)
                    scheme = sf.find('a:schemeClr', ns)
                    if srgb is not None:
                        border["color"] = f"#{srgb.get('val')}"
                    elif scheme is not None:
                        resolved = _resolve_scheme_color(scheme.get('val'), theme_colors, color_mapping)
                        if resolved:
                            border["color"] = resolved
            if border:
                borders[side] = border
        if borders:
            props["borders"] = borders

        # Vertical alignment
        anchor = tc_pr.get('anchor')
        if anchor:
            props["anchor"] = anchor  # t/ctr/b

        # Margins
        margins = {}
        for attr, key in [('marL', 'left'), ('marR', 'right'), ('marT', 'top'), ('marB', 'bottom')]:
            v = tc_pr.get(attr)
            if v:
                margins[key] = round(int(v) / 6350)
        if margins:
            props["margins"] = margins

    # Merge
    grid_span = tc.get('gridSpan')
    row_span = tc.get('rowSpan')
    if grid_span and int(grid_span) > 1:
        props["gridSpan"] = int(grid_span)
    if row_span and int(row_span) > 1:
        props["rowSpan"] = int(row_span)

    # Skip merged-away cells
    if tc.get('hMerge') == '1' or tc.get('vMerge') == '1':
        props["merged"] = True

    # Text styles (from first run)
    tf = cell.text_frame
    if tf.paragraphs:
        para = tf.paragraphs[0]
        if para.alignment is not None:
            align_map = {1: "left", 2: "center", 3: "right"}
            a = align_map.get(int(para.alignment))
            if a:
                props["align"] = a
        for run in para.runs:
            if run.font.bold:
                props["bold"] = True
            if run.font.italic:
                props["italic"] = True
            if run.font.size:
                props["fontSize"] = int(run.font.size.pt)
            try:
                fc = extract_text_color(run, theme_colors, color_mapping)
                if fc:
                    props["fontColor"] = fc
            except:
                pass
            break  # first run only

    if props:
        props["text"] = text
        return props
    return text


def _parse_table_style(pptx_path, style_id, theme_colors, color_mapping):
    """Parse tableStyles.xml and resolve fills/borders for a given style."""
    import zipfile
    ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    try:
        with zipfile.ZipFile(str(pptx_path)) as z:
            xml = z.read('ppt/tableStyles.xml')
        from lxml import etree
        root = etree.fromstring(xml)
        # Find style by ID, or use default
        target_id = style_id or root.get('def')
        style_el = None
        for s in root.findall('a:tblStyle', ns):
            if s.get('styleId') == target_id:
                style_el = s
                break
        if style_el is None:
            return {}

        def resolve_fill(tc_style):
            """Extract fill color from tcStyle element."""
            fill_el = tc_style.find('a:fill/a:solidFill', ns)
            if fill_el is None:
                return None
            srgb = fill_el.find('a:srgbClr', ns)
            if srgb is not None:
                return f"#{srgb.get('val')}"
            scheme = fill_el.find('a:schemeClr', ns)
            if scheme is not None:
                base = _resolve_scheme_color(scheme.get('val'), theme_colors, color_mapping)
                if base:
                    tint = scheme.find('a:tint', ns)
                    shade = scheme.find('a:shade', ns)
                    if tint is not None:
                        return _apply_tint(base, int(tint.get('val')) / 100000)
                    if shade is not None:
                        return _apply_shade(base, int(shade.get('val')) / 100000)
                    return base
            return None

        def resolve_border_color(tc_bdr):
            """Extract border color from tcBdr."""
            if tc_bdr is None:
                return None
            # Check any border line for color
            for tag in ['a:left', 'a:right', 'a:top', 'a:bottom', 'a:insideH', 'a:insideV']:
                ln = tc_bdr.find(f'{tag}/a:ln/a:solidFill', ns)
                if ln is not None:
                    scheme = ln.find('a:schemeClr', ns)
                    if scheme is not None:
                        return _resolve_scheme_color(scheme.get('val'), theme_colors, color_mapping)
                    srgb = ln.find('a:srgbClr', ns)
                    if srgb is not None:
                        return f"#{srgb.get('val')}"
            return None

        def resolve_text_color(tc_txt):
            """Extract text color from tcTxStyle."""
            if tc_txt is None:
                return None
            scheme = tc_txt.find('a:schemeClr', ns)
            if scheme is not None:
                return _resolve_scheme_color(scheme.get('val'), theme_colors, color_mapping)
            return None

        result = {}
        for part_name in ['wholeTbl', 'firstRow', 'lastRow', 'firstCol', 'lastCol', 'band1H', 'band2H']:
            part = style_el.find(f'a:{part_name}', ns)
            if part is None:
                continue
            info = {}
            tc_style = part.find('a:tcStyle', ns)
            if tc_style is not None:
                f = resolve_fill(tc_style)
                if f:
                    info['fill'] = f
                bc = resolve_border_color(tc_style.find('a:tcBdr', ns))
                if bc:
                    info['borderColor'] = bc
            tc_txt = part.find('a:tcTxStyle', ns)
            if tc_txt is not None:
                tc = resolve_text_color(tc_txt)
                if tc:
                    info['fontColor'] = tc
                if tc_txt.get('b') == 'on':
                    info['bold'] = True
            if info:
                result[part_name] = info
        return result
    except Exception:
        return {}


def _apply_tint(hex_color, tint_factor):
    """Apply tint (mix with white) to a hex color."""
    r, g, b = int(hex_color[1:3], 16), int(hex_color[3:5], 16), int(hex_color[5:7], 16)
    r = int(r + (255 - r) * (1 - tint_factor))
    g = int(g + (255 - g) * (1 - tint_factor))
    b = int(b + (255 - b) * (1 - tint_factor))
    return f"#{r:02X}{g:02X}{b:02X}"


def _apply_shade(hex_color, shade_factor):
    """Apply shade (mix with black) to a hex color."""
    r, g, b = int(hex_color[1:3], 16), int(hex_color[3:5], 16), int(hex_color[5:7], 16)
    r = int(r * shade_factor)
    g = int(g * shade_factor)
    b = int(b * shade_factor)
    return f"#{r:02X}{g:02X}{b:02X}"


def _apply_style_to_cell(cell_val, style_info):
    """Apply table style info to a cell that lacks explicit properties."""
    if not style_info:
        return cell_val
    is_str = isinstance(cell_val, str)
    if is_str:
        needs_upgrade = any(k in style_info for k in ('fill', 'fontColor', 'bold'))
        if not needs_upgrade:
            return cell_val
        cell_val = {"text": cell_val}
    # Only apply if cell doesn't already have the property
    for key in ('fill', 'fontColor', 'bold'):
        if key in style_info and key not in cell_val:
            cell_val[key] = style_info[key]
    return cell_val


def extract_table_element(shape, theme_colors=None, color_mapping=None, pptx_path=None):
    """Extract table as element dict."""
    try:
        table = shape.table
        ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        tbl_elem = table._tbl

        elem = {
            "type": "table",
            "x": round(shape.left / 6350),
            "y": round(shape.top / 6350),
            "width": round(shape.width / 6350),
            "height": round(shape.height / 6350),
        }

        # Column widths
        elem["colWidths"] = [round(col.width / 6350) for col in table.columns]

        # Row heights
        elem["rowHeights"] = [round(row.height / 6350) for row in table.rows]

        # Table style properties
        tbl_pr = tbl_elem.find('a:tblPr', ns)
        if tbl_pr is not None:
            for attr in ['firstRow', 'lastRow', 'firstCol', 'lastCol', 'bandRow', 'bandCol']:
                if tbl_pr.get(attr) == '1':
                    elem[attr] = True
            style_id = tbl_pr.find('a:tableStyleId', ns)
            if style_id is not None and style_id.text:
                elem["tableStyleId"] = style_id.text

        # Headers (first row)
        elem["headers"] = [_extract_cell(c, theme_colors, color_mapping) for c in table.rows[0].cells]

        # Data rows
        elem["rows"] = [
            [_extract_cell(c, theme_colors, color_mapping) for c in row.cells]
            for row in list(table.rows)[1:]
        ]

        # Apply table style fills/colors to cells without explicit values
        if pptx_path:
            style_id = elem.get("tableStyleId")
            ts = _parse_table_style(pptx_path, style_id, theme_colors, color_mapping)
            if ts:
                whole = ts.get('wholeTbl', {})
                first_row = {**whole, **ts.get('firstRow', {})} if elem.get('firstRow') else whole
                band1 = {**whole, **ts.get('band1H', {})} if elem.get('bandRow') else whole
                band2 = whole

                # Apply to headers
                elem["headers"] = [_apply_style_to_cell(c, first_row) for c in elem["headers"]]

                # Apply to data rows
                for ri, row in enumerate(elem["rows"]):
                    style = band1 if ri % 2 == 0 else band2
                    elem["rows"][ri] = [_apply_style_to_cell(c, style) for c in row]

        return elem
    except Exception as e:
        print(f"Warning: Failed to extract table: {e}", file=sys.stderr)
        return None

def detect_layout(slide):
    """Detect slide layout type."""
    # Check for specific layouts
    for shape in slide.shapes:
        if not shape.is_placeholder:
            continue
        
        ph_type = shape.placeholder_format.type
        
        # Check for title slide (has subtitle)
        if ph_type == 3:  # SUBTITLE
            return "title"
        
        # Check for thank you slide (specific layout)
        if ph_type == 1 and shape.text and "thank" in shape.text.lower():
            return "thankyou"
    
    # Check for content vs title_only
    has_title = False
    has_content = False
    
    for shape in slide.shapes:
        if shape.is_placeholder:
            ph_type = shape.placeholder_format.type
            if ph_type == 1:  # TITLE
                has_title = True
            elif ph_type in (2, 7, 13):  # BODY, OBJECT, CONTENT
                has_content = True
    
    if has_title and has_content:
        return "content"
    elif has_title:
        return "title_only"
    else:
        return "title_only"  # default

def extract_slide(slide, theme_colors=None, color_mapping=None, theme_styles=None, master_idx=0, pptx_path=None, slide_idx=0):
    """Extract slide content to dict."""
    slide_dict = {
        "layout": detect_layout(slide),
        "masterIndex": master_idx  # Add master index
    }
    
    # Extract title
    if slide.shapes.title:
        slide_dict["title"] = slide.shapes.title.text
    
    # Extract subtitle, date, etc for title layout
    if slide_dict["layout"] == "title":
        for shape in slide.shapes:
            if not shape.is_placeholder or not shape.text.strip():
                continue
            ph_type = shape.placeholder_format.type
            if ph_type == 4:  # SUBTITLE
                slide_dict["subtitle"] = shape.text
            elif ph_type == 2:  # BODY (used for date in title slide)
                if "date" not in slide_dict:
                    slide_dict["date"] = shape.text
    
    # Extract content from placeholders as textboxes (to preserve position)
    elements = []
    
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.type in (2, 7, 13):
            if shape.has_text_frame and shape.text.strip():
                # Create textbox element to preserve position
                textbox_elem = {
                    "type": "textbox",
                    "x": round(shape.left / 6350),
                    "y": round(shape.top / 6350),
                    "width": round(shape.width / 6350),
                    "height": round(shape.height / 6350),
                }
                
                # Check if it's a bullet list (multiple paragraphs)
                paragraphs_with_text = [p for p in shape.text_frame.paragraphs if p.text.strip()]
                
                if len(paragraphs_with_text) > 1:
                    # Multiple paragraphs - extract with bullet/numbering info
                    paragraphs = []
                    default_font_size = None
                    
                    for paragraph in paragraphs_with_text:
                        # Get font size
                        if not default_font_size:
                            for run in paragraph.runs:
                                if run.font.size:
                                    default_font_size = int(run.font.size.pt)
                                    break
                        # Check for bullet or numbering
                        has_bullet = False
                        numbering_type = None
                        bu_font = None
                        mar_l = None
                        indent = None
                        space_after = None
                        space_before = None
                        
                        try:
                            pPr = paragraph._element.pPr
                            
                            # Get from paragraph if available
                            if pPr is not None:
                                # Get margins
                                mar_l = pPr.get('marL')
                                indent = pPr.get('indent')
                                
                                # Get spacing
                                ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                                spc_aft = pPr.find('.//a:spcAft/a:spcPts', ns)
                                if spc_aft is not None:
                                    space_after = spc_aft.get('val')
                                
                                spc_bef = pPr.find('.//a:spcBef/a:spcPts', ns)
                                if spc_bef is not None:
                                    space_before = spc_bef.get('val')
                                
                                # Get bullet font
                                bu_font_elem = pPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}buFont')
                                if bu_font_elem is not None:
                                    bu_font = bu_font_elem.get('typeface')
                                
                                bu_auto_num = pPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}buAutoNum')
                                bu_char = pPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}buChar')
                                
                                if bu_auto_num is not None:
                                    numbering_type = bu_auto_num.get('type', 'arabicPeriod')
                                elif bu_char is not None:
                                    has_bullet = True
                            
                            # If no pPr or missing values, get defaults from layout placeholder
                            if pPr is None or not space_after or not mar_l or not bu_font or not numbering_type:
                                try:
                                    import zipfile
                                    import xml.etree.ElementTree as ET
                                    
                                    with zipfile.ZipFile(pptx_path, 'r') as zip_ref:
                                        layout_idx = list(slide.slide_layout.slide_master.slide_layouts).index(slide.slide_layout)
                                        layout_xml = zip_ref.read(f'ppt/slideLayouts/slideLayout{layout_idx + 1}.xml')
                                        root = ET.fromstring(layout_xml)
                                        
                                        ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
                                              'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                                        
                                        # Find placeholder with matching type
                                        ph_type_map = {2: 'body', 7: 'obj', 13: 'body'}
                                        ph_type_str = ph_type_map.get(shape.placeholder_format.type, 'body')
                                        
                                        for sp in root.findall('.//p:sp', ns):
                                            ph = sp.find(f'.//p:ph[@type="{ph_type_str}"]', ns)
                                            if ph is not None:
                                                # Get lstStyle spacing
                                                lst_style = sp.find('.//a:lstStyle/a:lvl1pPr', ns)
                                                if lst_style is not None:
                                                    if not space_after:
                                                        spc_aft = lst_style.find('.//a:spcAft/a:spcPts', ns)
                                                        if spc_aft is not None:
                                                            space_after = spc_aft.get('val')
                                                    
                                                    # Get margins from layout if not in paragraph
                                                    if not mar_l:
                                                        mar_l = lst_style.get('marL')
                                                    if not indent:
                                                        indent = lst_style.get('indent')
                                                    
                                                    # Get buFont
                                                    if not bu_font:
                                                        bu_font_elem = lst_style.find('.//a:buFont', ns)
                                                        if bu_font_elem is not None:
                                                            bu_font = bu_font_elem.get('typeface')
                                                    
                                                    # Check for numbering
                                                    if not numbering_type:
                                                        bu_auto_num = lst_style.find('.//a:buAutoNum', ns)
                                                        if bu_auto_num is not None:
                                                            numbering_type = bu_auto_num.get('type', 'arabicPeriod')
                                                break
                                except:
                                    pass
                        except:
                            pass
                        
                        parts = []
                        for run in paragraph.runs:
                            if not run.text:
                                continue
                            
                            # Check for hyperlink
                            if run.hyperlink and run.hyperlink.address:
                                parts.append(f"{{{{link:{run.hyperlink.address}:{run.text}}}}}")
                            else:
                                styles = []
                                if run.font.bold:
                                    styles.append("bold")
                                if run.font.italic:
                                    styles.append("italic")
                                if run.font.size:
                                    pt = int(run.font.size.pt)
                                    if pt != 18:
                                        styles.append(f"{pt}pt")
                                
                                try:
                                    hex_color = extract_text_color(run, theme_colors, color_mapping, is_placeholder=True)
                                    if hex_color and hex_color != "#000000":
                                        styles.append(hex_color)
                                except:
                                    pass
                                
                                if styles:
                                    parts.append(f"{{{{{','.join(styles)}:{run.text}}}}}")
                                else:
                                    parts.append(run.text)
                        
                        item_text = ''.join(parts)
                        if item_text.strip():
                            para_info = {"text": item_text}
                            if has_bullet:
                                para_info["bullet"] = True
                            elif numbering_type:
                                para_info["numbering"] = numbering_type
                            
                            # Add bullet font
                            if bu_font:
                                para_info["buFont"] = bu_font
                            
                            # Add margin/indent if specified
                            if mar_l:
                                para_info["marL"] = int(mar_l)
                            if indent:
                                para_info["indent"] = int(indent)
                            if space_after:
                                para_info["spaceAfter"] = int(space_after)
                            if space_before:
                                para_info["spaceBefore"] = int(space_before)
                            
                            paragraphs.append(para_info)
                    
                    textbox_elem["paragraphs"] = paragraphs
                    
                    # Add fontSize - try to get from master txStyles
                    if not default_font_size:
                        try:
                            import zipfile
                            import xml.etree.ElementTree as ET
                            
                            with zipfile.ZipFile(pptx_path, 'r') as zip_ref:
                                # Use master_idx from function parameter
                                master_xml = zip_ref.read(f'ppt/slideMasters/slideMaster{master_idx + 1}.xml')
                                root = ET.fromstring(master_xml)
                                
                                ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
                                      'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                                
                                # Get from bodyStyle
                                body_style = root.find('.//p:txStyles/p:bodyStyle/a:lvl1pPr/a:defRPr[@sz]', ns)
                                if body_style is not None:
                                    sz = body_style.get('sz')
                                    if sz:
                                        default_font_size = int(sz) // 100
                        except:
                            pass
                    
                    if default_font_size and default_font_size != 18:
                        textbox_elem["fontSize"] = default_font_size
                    
                    elements.append(textbox_elem)
                else:
                    # Single paragraph - extract as textbox to preserve position
                    paragraph = paragraphs_with_text[0]
                    textbox_elem = {
                        "type": "textbox",
                        "x": round(shape.left / 6350),
                        "y": round(shape.top / 6350),
                        "width": round(shape.width / 6350),
                    }
                    
                    parts = []
                    for run in paragraph.runs:
                        text = run.text
                        if not text:
                            continue
                        
                        # Check for hyperlink
                        if run.hyperlink and run.hyperlink.address:
                            parts.append(f"{{{{link:{run.hyperlink.address}:{text}}}}}")
                        else:
                            styles = []
                            if run.font.bold:
                                styles.append("bold")
                            if run.font.italic:
                                styles.append("italic")
                            if run.font.size:
                                pt = int(run.font.size.pt)
                                if pt != 18:
                                    styles.append(f"{pt}pt")
                            
                            try:
                                hex_color = extract_text_color(run, theme_colors, color_mapping, is_placeholder=True)
                                if hex_color and hex_color != "#000000":
                                    styles.append(hex_color)
                            except:
                                pass
                            
                            if styles:
                                parts.append(f"{{{{{','.join(styles)}:{text}}}}}")
                            else:
                                parts.append(text)
                    
                    item_text = ''.join(parts)
                    if item_text.strip():
                        textbox_elem["text"] = item_text
                        
                        # Get font size
                        for run in paragraph.runs:
                            if run.font.size:
                                textbox_elem["fontSize"] = int(run.font.size.pt)
                                break
                        
                        # Get alignment
                        if paragraph.alignment == 2:
                            textbox_elem["align"] = "center"
                        elif paragraph.alignment == 3:
                            textbox_elem["align"] = "right"
                        
                        elements.append(textbox_elem)
    
    # Extract non-placeholder elements (elements already initialized above)
    img_counter = 0
    for shape in slide.shapes:
        if shape.is_placeholder:
            continue
        
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                output_dir = pptx_path.parent if pptx_path else None
                elem, img_counter = extract_group_element(shape, theme_colors, color_mapping, theme_styles, output_dir, slide_idx, img_counter)
                if elem:
                    elements.append(elem)
            elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                elem = extract_textbox_element(shape, theme_colors, color_mapping, theme_styles)
                if elem:
                    elements.append(elem)
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                output_dir = pptx_path.parent if pptx_path else None
                elem = extract_picture_element(shape, output_dir, slide_idx, img_counter)
                if elem:
                    elements.append(elem)
                    img_counter += 1
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                elem = extract_table_element(shape, theme_colors, color_mapping, pptx_path)
                if elem:
                    elements.append(elem)
            elif shape.shape_type == MSO_SHAPE_TYPE.LINE:
                elem = extract_line_element(shape, theme_colors, color_mapping)
                if elem:
                    elements.append(elem)
            elif shape.shape_type in (MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM):
                # Try freeform path extraction first
                if shape.shape_type == MSO_SHAPE_TYPE.FREEFORM:
                    elem = extract_freeform_element(shape, theme_colors, color_mapping)
                    if elem:
                        elements.append(elem)
                        continue
                elem = extract_shape_element(shape, theme_colors, color_mapping, theme_styles)
                if elem:
                    elements.append(elem)
        except Exception as e:
            print(f"Warning: Failed to extract shape {shape.name}: {e}", file=sys.stderr)
    
    if elements:
        slide_dict["elements"] = elements
    
    return slide_dict

def pptx_to_json(pptx_path: Path, output_path: Path = None):
    """Convert PPTX to JSON."""
    prs = Presentation(str(pptx_path))
    
    # Detect theme (simple heuristic)
    theme = "light"  # default
    
    result = {
        "theme": theme,
        "slides": []
    }
    
    for slide_idx, slide in enumerate(prs.slides):
        # Get slide master index
        slide_master = slide.slide_layout.slide_master
        master_idx = list(prs.slide_masters).index(slide_master)
        
        # Extract theme colors and mapping for this master
        theme_colors, color_mapping, theme_styles = extract_theme_colors_and_mapping(pptx_path, master_idx)
        
        slide_dict = extract_slide(slide, theme_colors, color_mapping, theme_styles, master_idx, pptx_path, slide_idx)
        result["slides"].append(slide_dict)
    
    # Output
    if output_path:
        output_path.write_text(json.dumps(result, indent=2, ensure_ascii=False))
        print(f"Converted: {output_path}")
    else:
        print(json.dumps(result, indent=2, ensure_ascii=False))
    
    return result

def main():
    parser = argparse.ArgumentParser(description="Convert PPTX to JSON")
    parser.add_argument("input", help="Input PPTX file")
    parser.add_argument("-o", "--output", help="Output JSON file (default: stdout)")
    args = parser.parse_args()
    
    input_path = Path(args.input)
    if not input_path.exists():
        print(f"Error: File not found: {input_path}", file=sys.stderr)
        sys.exit(1)
    
    output_path = Path(args.output) if args.output else None
    pptx_to_json(input_path, output_path)

if __name__ == "__main__":
    main()

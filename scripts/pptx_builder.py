#!/usr/bin/env python3
"""PPTX Builder - Generate PowerPoint from JSON using template"""
import argparse
import json
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_CONNECTOR

# Icon settings
ICON_DIR = Path.home() / "Library/Application Support/com.raycast.macos/extensions/aws-icons/icons"
ICON_LOCAL_DIR = Path(__file__).parent.parent / "icons"  # Fallback: pptx-maker/icons/

# Layout mapping for each theme (master index -> layout index)
LAYOUT_MAP = {
    "light": {  # Master 0
        "title": 0,        # Title Side 1A
        "agenda": 7,       # Agenda Slide 2
        "section": 39,     # Section Header Option 1
        "subsection": 38,  # Section Header Option 2
        "content": 7,      # 通常スライド（agendaの見た目＋自由配置）
        "title_only": 9,   # タイトルのみ（広く使いたい時）
        "thankyou": 47,    # Thank You Option 1
    },
    "dark": {  # Master 1
        "title": 0,        # Title Side 1A
        "agenda": 7,       # Agenda Slide 2
        "section": 39,     # Section Header Option 1
        "subsection": 38,  # Section Header Option 2
        "content": 9,      # Title Only with Left Line
        "title_only": 10,  # Title Only
        "thankyou": 47,    # Thank You Option 1
    }
}

# Brand fonts
FONT_FULLWIDTH = "メイリオ"
FONT_HALFWIDTH = "Amazon Ember"

# Theme color palette
THEME_COLORS = {
    "light": {
        "text": RGBColor(0x00, 0x00, 0x00),          # Black
        "text_inverse": RGBColor(0xFF, 0xFF, 0xFF),  # White
        "background": RGBColor(0xFF, 0xFF, 0xFF),    # White
        "background_alt": RGBColor(0xF3, 0xF3, 0xF3),# Light gray
        "accent": RGBColor(0xFF, 0x99, 0x00),        # AWS Orange
        "header": RGBColor(0x5F, 0x63, 0x68),        # Gray
    },
    "dark": {
        "text": RGBColor(0xFF, 0xFF, 0xFF),          # White
        "text_inverse": RGBColor(0x00, 0x00, 0x00),  # Black
        "background": RGBColor(0x23, 0x2F, 0x3E),    # AWS Dark
        "background_alt": RGBColor(0x2E, 0x3B, 0x4E),# Slightly lighter
        "accent": RGBColor(0xFF, 0x99, 0x00),        # AWS Orange
        "header": RGBColor(0xFF, 0x99, 0x00),        # AWS Orange
    }
}

def is_fullwidth(char):
    """全角文字判定"""
    return ord(char) > 127

def normalize_spacing(text):
    """半角と全角の間に半角スペースを挿入（括弧内側は除外）"""
    # 開き括弧
    open_br = set('([{「『（【〔《〈')
    # 閉じ括弧
    close_br = set(')]}」』）】〕》〉')
    # ペア追跡が必要な引用符
    quotes = set('"\'`')
    
    def is_ascii(c):
        return '\x21' <= c <= '\x7e'
    
    def is_wide(c):
        return c > '\x7f' and c not in open_br and c not in close_br
    
    # 引用符の開閉状態を追跡（先読み用）
    quote_positions = {q: [] for q in quotes}
    for i, ch in enumerate(text):
        if ch in quotes:
            quote_positions[ch].append(i)
    
    def is_opening_quote(pos, ch):
        if ch not in quotes:
            return False
        positions = quote_positions[ch]
        idx = positions.index(pos)
        return idx % 2 == 0
    
    def is_closing_quote(pos, ch):
        if ch not in quotes:
            return False
        positions = quote_positions[ch]
        idx = positions.index(pos)
        return idx % 2 == 1
    
    result = []
    for i, ch in enumerate(text):
        result.append(ch)
        if i + 1 >= len(text):
            continue
        
        next_ch = text[i + 1]
        
        # 閉じ引用符の後はスペース（次が開き括弧・開き引用符以外）
        if is_closing_quote(i, ch) and next_ch not in open_br and not is_opening_quote(i + 1, next_ch):
            result.append(' ')
        # 開き引用符の前はスペース
        elif is_opening_quote(i + 1, next_ch) and ch not in open_br and not is_opening_quote(i, ch):
            result.append(' ')
        # 半角→全角の境界（開き括弧・開き引用符直後は除外）
        elif is_ascii(ch) and is_wide(next_ch) and ch not in open_br and not is_opening_quote(i, ch):
            result.append(' ')
        # 全角→半角の境界（閉じ括弧・閉じ引用符直前は除外）
        elif is_wide(ch) and is_ascii(next_ch) and next_ch not in close_br and not is_closing_quote(i + 1, next_ch):
            result.append(' ')
    
    return ''.join(result)


def parse_styled_text(text):
    """Parse {{attrs:text}} syntax into styled segments.
    
    Supported attrs (comma-separated):
    - bold
    - italic  
    - #RRGGBB (color)
    - NNpt (font size)
    - link:URL (hyperlink)
    
    Examples:
    - {{bold:太字}}
    - {{#FF9900:オレンジ}}
    - {{bold,#FF9900:太字オレンジ}}
    - {{24pt:大きい文字}}
    - {{bold,18pt,#FF9900:全部入り}}
    - {{link:https://aws.amazon.com:AWS公式サイト}}
    """
    # First, handle link syntax specially: {{link:URL:text}}
    link_pattern = r'\{\{link:([^}]+)\}\}'
    segments = []
    last_end = 0
    
    for match in re.finditer(link_pattern, text):
        # Add plain text before this match
        if match.start() > last_end:
            plain_part = text[last_end:match.start()]
            segments.extend(_parse_non_link_styles(plain_part))
        
        # Split by last colon to separate URL and text
        content = match.group(1)
        last_colon = content.rfind(':')
        if last_colon > 0:
            url = content[:last_colon]
            link_text = content[last_colon+1:]
            segments.append({"text": link_text, "link": url})
        else:
            # No text specified, use URL as text
            segments.append({"text": content, "link": content})
        
        last_end = match.end()
    
    # Process remaining text
    if last_end < len(text):
        segments.extend(_parse_non_link_styles(text[last_end:]))
    
    return segments if segments else [{"text": text}]

def _parse_non_link_styles(text):
    """Parse non-link styled text."""
    pattern = r'\{\{([^:}]+):([^}]+)\}\}'
    segments = []
    last_end = 0
    
    for match in re.finditer(pattern, text):
        # Add plain text before this match
        if match.start() > last_end:
            plain_text = normalize_spacing(text[last_end:match.start()])
            segments.append({"text": plain_text})
        
        # Parse attributes
        attrs = match.group(1).split(',')
        inner_text = normalize_spacing(match.group(2))
        segment = {"text": inner_text}
        
        for attr in attrs:
            attr = attr.strip()
            if attr == "bold":
                segment["bold"] = True
            elif attr == "italic":
                segment["italic"] = True
            elif attr.startswith("#") and len(attr) == 7:
                segment["color"] = attr
            elif attr.endswith("pt"):
                try:
                    segment["fontSize"] = int(attr[:-2])
                except ValueError:
                    pass
        
        segments.append(segment)
        last_end = match.end()
    
    # Add remaining plain text
    if last_end < len(text):
        plain_text = normalize_spacing(text[last_end:])
        segments.append({"text": plain_text})
    
    # Insert spacing at segment boundaries (half-width ↔ full-width)
    for i in range(len(segments) - 1):
        cur_text = segments[i]["text"]
        nxt_text = segments[i + 1]["text"]
        if not cur_text or not nxt_text:
            continue
        last_ch, first_ch = cur_text[-1], nxt_text[0]
        is_ascii = lambda c: '\x21' <= c <= '\x7e'
        is_wide = lambda c: c > '\x7f'
        if (is_ascii(last_ch) and is_wide(first_ch)) or (is_wide(last_ch) and is_ascii(first_ch)):
            segments[i]["text"] = cur_text + ' '
    
    return segments


def _expand_styled_newlines(text):
    """Expand \\n inside styled tags so each line gets its own complete tag.
    
    e.g. '{{bold,#FF9900:line1\\nline2}}' -> '{{bold,#FF9900:line1}}\\n{{bold,#FF9900:line2}}'
    """
    # Match styled tags that contain \n in their content
    # Handle both link and non-link styles
    def expand_match(m):
        attrs = m.group(1)
        content = m.group(2)
        if '\n' not in content:
            return m.group(0)
        lines = content.split('\n')
        return '\n'.join(f'{{{{{attrs}:{line}}}}}' for line in lines)
    
    # Non-link: {{attrs:content}}
    text = re.sub(r'\{\{([^:}]+):([^}]*\n[^}]*)\}\}', expand_match, text)
    return text


def _icons_not_installed_error():
    """Print icon installation instructions and exit."""
    print("=" * 60, file=sys.stderr)
    print("CRITICAL: Icons not installed. Cannot continue.", file=sys.stderr)
    print("=" * 60, file=sys.stderr)
    print("", file=sys.stderr)
    print("Icons are required for slide generation.", file=sys.stderr)
    if sys.platform == "darwin":
        print("", file=sys.stderr)
        print("  [Recommended] Install Raycast AWS Icons extension:", file=sys.stderr)
        print("    https://gitlab.aws.dev/sktok/Raycast-AWS-Icons", file=sys.stderr)
        print("", file=sys.stderr)
        print("  [Alternative] Download icons directly:", file=sys.stderr)
        print("    uv run python3 scripts/download_icons.py", file=sys.stderr)
    else:
        print("", file=sys.stderr)
        print("  Run: uv run python3 scripts/download_icons.py", file=sys.stderr)
    print("", file=sys.stderr)
    print("Stop current work and ask the user which option to use.", file=sys.stderr)
    sys.exit(1)


def check_icon_exists(name: str, theme: str = "light") -> bool:
    """Check if icon exists without raising error."""
    search_dirs = []
    if ICON_DIR.exists():
        search_dirs.append(ICON_DIR)
    if ICON_LOCAL_DIR.exists():
        search_dirs.append(ICON_LOCAL_DIR)
    
    if not search_dirs:
        return False
    
    for dir_path in search_dirs:
        for ext in [".svg", ".png", ".gif", ".jpg", ".jpeg"]:
            path = dir_path / f"{name}{ext}"
            if path.exists():
                return True
    return False


def resolve_icon_path(name: str, theme: str = "light") -> Path:
    """Resolve icon name to file path with fallback.
    
    Search order:
    1. Raycast extension icons (ICON_DIR)
    2. Local pptx-maker/icons/ (ICON_LOCAL_DIR)
    """
    search_dirs = []
    if ICON_DIR.exists():
        search_dirs.append(ICON_DIR)
    if ICON_LOCAL_DIR.exists():
        search_dirs.append(ICON_LOCAL_DIR)
    
    if not search_dirs:
        _icons_not_installed_error()
    
    for dir_path in search_dirs:
        for ext in [".svg", ".png", ".gif", ".jpg", ".jpeg"]:
            path = dir_path / f"{name}{ext}"
            if path.exists():
                return path
    
    print(f"Error: Icon not found: {name}", file=sys.stderr)
    print(f"  Searched: {', '.join(str(d) for d in search_dirs)}", file=sys.stderr)
    if not ICON_LOCAL_DIR.exists():
        print(f"", file=sys.stderr)
        print(f"Icons may be incomplete. Try:", file=sys.stderr)
        print(f"  uv run python3 scripts/download_icons.py", file=sys.stderr)
    raise FileNotFoundError(f"Icon not found: {name}")


def resolve_image_path(src: str, theme: str = "light") -> Path:
    """Resolve image source to file path."""
    if src.startswith("icons:"):
        return resolve_icon_path(src[6:], theme)
    else:
        path = Path(src).expanduser()
        if not path.exists():
            raise FileNotFoundError(f"Image not found: {src}")
        return path


def _recolor_svg(svg_bytes: bytes, color: str) -> bytes | None:
    """Recolor single-color SVG. Returns None if multi-color (skip)."""
    import re
    text = svg_bytes.decode("utf-8")
    
    # Collect real colors from fill/stroke (attributes + CSS)
    attr_colors = re.findall(r'(?:fill|stroke)\s*[=:]\s*["\']?\s*(#[0-9a-fA-F]{3,8}|rgb[^)]*\))', text)
    unique = set(c.lower().strip() for c in attr_colors)
    
    if len(unique) == 0:
        return None
    if len(unique) > 1:
        print(f"Warning: iconColor skipped (multi-color SVG, {len(unique)} colors found)", file=sys.stderr)
        return None
    
    original = unique.pop()
    
    # Determine if fill-based or stroke-based
    has_fill = bool(re.search(r'fill\s*[=:]\s*["\']?\s*' + re.escape(original), text, re.IGNORECASE))
    has_stroke = bool(re.search(r'stroke\s*[=:]\s*["\']?\s*' + re.escape(original), text, re.IGNORECASE))
    
    if has_fill:
        text = re.sub(r'(fill\s*[=:]\s*["\']?\s*)' + re.escape(original), lambda m: m.group(1) + color, text, flags=re.IGNORECASE)
    if has_stroke:
        text = re.sub(r'(stroke\s*[=:]\s*["\']?\s*)' + re.escape(original), lambda m: m.group(1) + color, text, flags=re.IGNORECASE)
    
    return text.encode("utf-8")


def get_svg_dimensions(svg_path: Path) -> tuple[int, int]:
    """Get SVG dimensions from viewBox or width/height attributes."""
    from lxml import etree
    tree = etree.parse(str(svg_path))
    root = tree.getroot()
    ns = {'svg': 'http://www.w3.org/2000/svg'}
    
    # Try viewBox first
    vb = root.get('viewBox')
    if vb:
        parts = vb.replace(',', ' ').split()
        if len(parts) == 4:
            return int(float(parts[2])), int(float(parts[3]))
    
    # Fallback to width/height
    w = root.get('width', '').replace('px', '')
    h = root.get('height', '').replace('px', '')
    if w and h:
        try:
            return int(float(w)), int(float(h))
        except ValueError:
            pass
    return 100, 100  # Fallback square


def add_svg_to_slide(slide, svg_bytes: bytes, x, y, width, height):
    """Add SVG directly to PPTX slide via OpenXML asvg:svgBlip."""
    from lxml import etree
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    from pptx.opc.package import Part
    from pptx.opc.packuri import PackURI

    slide_part = slide.part

    # Find unique partname
    idx = 1
    while True:
        partname = f'/ppt/media/svg_image{idx}.svg'
        if not any(partname == str(p.partname) for p in slide_part.package.iter_parts()):
            break
        idx += 1

    svg_part = Part(PackURI(partname), 'image/svg+xml', slide_part.package, svg_bytes)
    rId = slide_part.relate_to(svg_part, RT.IMAGE)

    spTree = slide.shapes._spTree
    shape_id = max((int(sp.get('id', 0)) for sp in spTree.iter()), default=0) + 1

    SVG_BLIP_URI = '{96DAC541-7B7A-43D3-8B79-37D633B846F1}'
    pic_xml = (
        f'<p:pic xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        f' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
        f' xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"'
        f' xmlns:asvg="http://schemas.microsoft.com/office/drawing/2016/SVG/main">'
        f'<p:nvPicPr>'
        f'<p:cNvPr id="{shape_id}" name="SVG {shape_id}"/>'
        f'<p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>'
        f'<p:nvPr/>'
        f'</p:nvPicPr>'
        f'<p:blipFill>'
        f'<a:blip r:embed="{rId}">'
        f'<a:extLst>'
        f'<a:ext uri="{SVG_BLIP_URI}">'
        f'<asvg:svgBlip r:embed="{rId}"/>'
        f'</a:ext>'
        f'</a:extLst>'
        f'</a:blip>'
        f'<a:stretch><a:fillRect/></a:stretch>'
        f'</p:blipFill>'
        f'<p:spPr>'
        f'<a:xfrm>'
        f'<a:off x="{x}" y="{y}"/>'
        f'<a:ext cx="{width}" cy="{height}"/>'
        f'</a:xfrm>'
        f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
        f'</p:spPr>'
        f'</p:pic>'
    )

    pic_element = etree.fromstring(pic_xml)
    spTree.append(pic_element)
    return pic_element


class PPTXBuilder:
    # Slide dimensions in EMU (16:9)
    SLIDE_WIDTH = 12192000
    SLIDE_HEIGHT = 6858000
    # px to EMU conversion (1920x1080 basis)
    EMU_PER_PX = 6350  # 12192000 / 1920 = 6858000 / 1080
    
    def __init__(self, template_path: Path, theme: str = "light"):
        self.prs = Presentation(str(template_path))
        self.theme = theme
        self.master_idx = 0 if theme == "light" else 1
        self.layouts = LAYOUT_MAP[theme]
        self.colors = THEME_COLORS[theme]
        self._clear_slides()
    
    def _clear_slides(self):
        while len(self.prs.slides) > 0:
            rId = self.prs.slides._sldIdLst[0].rId
            self.prs.part.drop_rel(rId)
            del self.prs.slides._sldIdLst[0]
    
    def _get_layout(self, layout_name: str):
        if layout_name not in self.layouts:
            raise ValueError(f"Unknown layout: {layout_name}. Available: {list(self.layouts.keys())}")
        layout_idx = self.layouts[layout_name]
        return self.prs.slide_masters[self.master_idx].slide_layouts[layout_idx]
    
    def add_slide(self, slide_def: dict):
        layout_name = slide_def.get("layout", "content")
        master_idx = slide_def.get("masterIndex", self.master_idx)  # Use slide's master or default
        
        # Get layout from appropriate master
        layout_idx = self.layouts[layout_name]
        layout = self.prs.slide_masters[master_idx].slide_layouts[layout_idx]
        slide = self.prs.slides.add_slide(layout)
        
        if layout_name == "title":
            self._fill_title(slide, slide_def)
        elif layout_name == "agenda":
            self._fill_content(slide, slide_def)
        elif layout_name == "content":
            self._fill_title_only(slide, slide_def)
        elif layout_name in ("section", "subsection"):
            self._fill_section(slide, slide_def)
        elif layout_name == "title_only":
            self._fill_title_only(slide, slide_def)
        # thankyou: no content to fill
        
        # Process elements for any layout
        for elem in slide_def.get("elements", []):
            elem_type = elem.get("type")
            if elem_type == "group":
                self._add_group(slide, elem)
            elif elem_type == "table":
                self._add_table(slide, elem)
            elif elem_type == "textbox":
                self._add_textbox(slide, elem)
            elif elem_type == "image":
                self._add_image(slide, elem)
            elif elem_type == "shape":
                self._add_shape(slide, elem)
            elif elem_type == "freeform":
                self._add_freeform_shape(slide, elem)
            elif elem_type == "line":
                self._add_line(slide, elem)
        
        # Add speaker notes
        if "notes" in slide_def:
            notes_frame = slide.notes_slide.notes_text_frame
            notes_frame.clear()
            p = notes_frame.paragraphs[0]
            self._apply_styled_text(p, slide_def["notes"])
        
        return slide
    
    def _fill_title(self, slide, d):
        if "title" in d and slide.shapes.title:
            slide.shapes.title.text = normalize_spacing(d["title"])
        for ph in slide.placeholders:
            idx = ph.placeholder_format.idx
            if idx == 1 and "subtitle" in d:
                ph.text = normalize_spacing(d["subtitle"])
            elif idx == 11 and "department" in d:
                ph.text = normalize_spacing(d["department"])
            elif idx == 12 and "date" in d:
                ph.text = normalize_spacing(d["date"])
            elif idx == 13 and "customer" in d:
                ph.text = normalize_spacing(d["customer"])
    
    def _fill_content(self, slide, d):
        if "title" in d and slide.shapes.title:
            slide.shapes.title.text = normalize_spacing(d["title"])
        for ph in slide.placeholders:
            if ph.placeholder_format.idx in (1, 10, 13):  # BODY/OBJECT placeholder
                if "items" in d:
                    tf = ph.text_frame
                    for i, item in enumerate(d["items"]):
                        if i == 0:
                            p = tf.paragraphs[0]
                        else:
                            p = tf.add_paragraph()
                        p.level = 0
                        self._set_bullet(p)
                        self._apply_styled_text(p, item)
                elif "body" in d:
                    self._apply_styled_text(ph.text_frame.paragraphs[0], d["body"])
                break
    
    def _set_bullet(self, paragraph):
        """Set bullet point formatting on paragraph."""
        from pptx.oxml.ns import qn
        from lxml import etree
        
        pPr = paragraph._element.get_or_add_pPr()
        pPr.set('marL', '285750')
        pPr.set('indent', '-285750')
        
        # Add bullet font
        buFont = etree.SubElement(pPr, qn('a:buFont'))
        buFont.set('typeface', 'Arial')
        
        # Add bullet character
        buChar = etree.SubElement(pPr, qn('a:buChar'))
        buChar.set('char', '•')
    
    def _set_numbering(self, paragraph, numbering_type='arabicPeriod'):
        """Set numbering formatting on paragraph."""
        from pptx.oxml.ns import qn
        from lxml import etree
        
        pPr = paragraph._element.get_or_add_pPr()
        pPr.set('marL', '285750')
        pPr.set('indent', '-285750')
        
        # Add auto numbering
        buAutoNum = etree.SubElement(pPr, qn('a:buAutoNum'))
        buAutoNum.set('type', numbering_type)
    
    def _fill_section(self, slide, d):
        if "title" in d and slide.shapes.title:
            slide.shapes.title.text = normalize_spacing(d["title"])
    
    def _fill_title_only(self, slide, d):
        self._set_or_remove_title(slide, d)

    def _set_or_remove_title(self, slide, d):
        """Set title text or remove title placeholder if empty."""
        if "title" not in d:
            return
        if d["title"] and slide.shapes.title:
            slide.shapes.title.text = normalize_spacing(d["title"])
        elif not d["title"] and slide.shapes.title:
            sp = slide.shapes.title._element
            sp.getparent().remove(sp)

    def _add_group(self, slide, elem):
        """Add group element to slide (flatten sub-elements)."""
        for sub_elem in elem.get("elements", []):
            sub_type = sub_elem.get("type")
            if sub_type == "group":
                self._add_group(slide, sub_elem)
            elif sub_type == "textbox":
                self._add_textbox(slide, sub_elem)
            elif sub_type == "image":
                self._add_image(slide, sub_elem)
            elif sub_type == "shape":
                self._add_shape(slide, sub_elem)
            elif sub_type == "freeform":
                self._add_freeform_shape(slide, sub_elem)
            elif sub_type == "line":
                self._add_line(slide, sub_elem)

    def _add_freeform_shape(self, slide, elem):
        """Add freeform shape with custom geometry."""
        from lxml import etree
        from pptx.oxml.ns import qn
        
        x_emu = self._px_to_emu(elem.get("x", 0))
        y_emu = self._px_to_emu(elem.get("y", 0))
        width_emu = self._px_to_emu(elem.get("width", 10))
        height_emu = self._px_to_emu(elem.get("height", 10))
        
        # Create rectangle as placeholder
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_emu, y_emu, width_emu, height_emu)
        
        # Build custom geometry from path array or raw XML
        path_cmds = elem.get("path")
        custom_geom = elem.get("customGeometry")
        
        if path_cmds:
            try:
                ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
                cust_geom = etree.SubElement(etree.Element('dummy'), f'{{{ns_a}}}custGeom')
                etree.SubElement(cust_geom, f'{{{ns_a}}}avLst')
                etree.SubElement(cust_geom, f'{{{ns_a}}}gdLst')
                etree.SubElement(cust_geom, f'{{{ns_a}}}ahLst')
                etree.SubElement(cust_geom, f'{{{ns_a}}}cxnLst')
                rect = etree.SubElement(cust_geom, f'{{{ns_a}}}rect')
                rect.set('l', 'l'); rect.set('t', 't'); rect.set('r', 'r'); rect.set('b', 'b')
                path_lst = etree.SubElement(cust_geom, f'{{{ns_a}}}pathLst')
                path_el = etree.SubElement(path_lst, f'{{{ns_a}}}path')
                path_el.set('w', str(width_emu))
                path_el.set('h', str(height_emu))
                
                for cmd in path_cmds:
                    c = cmd.get("cmd", "")
                    if c == "M":
                        mv = etree.SubElement(path_el, f'{{{ns_a}}}moveTo')
                        pt = etree.SubElement(mv, f'{{{ns_a}}}pt')
                        pt.set('x', str(round(cmd["x"] * self.EMU_PER_PX)))
                        pt.set('y', str(round(cmd["y"] * self.EMU_PER_PX)))
                    elif c == "L":
                        ln = etree.SubElement(path_el, f'{{{ns_a}}}lnTo')
                        pt = etree.SubElement(ln, f'{{{ns_a}}}pt')
                        pt.set('x', str(round(cmd["x"] * self.EMU_PER_PX)))
                        pt.set('y', str(round(cmd["y"] * self.EMU_PER_PX)))
                    elif c == "C":
                        cb = etree.SubElement(path_el, f'{{{ns_a}}}cubicBezTo')
                        for p in cmd["pts"]:
                            pt = etree.SubElement(cb, f'{{{ns_a}}}pt')
                            pt.set('x', str(round(p[0] * self.EMU_PER_PX)))
                            pt.set('y', str(round(p[1] * self.EMU_PER_PX)))
                    elif c == "Z":
                        etree.SubElement(path_el, f'{{{ns_a}}}close')
                
                sp_pr = shape._element.spPr
                for prst_geom in sp_pr.findall(qn('a:prstGeom')):
                    sp_pr.remove(prst_geom)
                xfrm = sp_pr.find(qn('a:xfrm'))
                if xfrm is not None:
                    xfrm.addnext(cust_geom)
                
                parent = shape._element
                for style in parent.findall(qn('p:style')):
                    parent.remove(style)
            except Exception as e:
                print(f"Warning: Failed to build freeform path: {e}", file=sys.stderr)
        elif custom_geom:
            try:
                custom_geom_clean = custom_geom.replace('ns0:', 'a:').replace('xmlns:ns0=', 'xmlns:a=')
                cust_geom_elem = etree.fromstring(custom_geom_clean)
                sp_pr = shape._element.spPr
                
                for prst_geom in sp_pr.findall(qn('a:prstGeom')):
                    sp_pr.remove(prst_geom)
                
                xfrm = sp_pr.find(qn('a:xfrm'))
                if xfrm is not None:
                    children = list(sp_pr)
                    xfrm_idx = children.index(xfrm)
                    sp_pr.insert(xfrm_idx + 1, cust_geom_elem)
                
                # Remove style element
                parent = shape._element
                for style in parent.findall(qn('p:style')):
                    parent.remove(style)
            except Exception as e:
                print(f"Warning: Failed to apply freeform geometry: {e}", file=sys.stderr)
        
        # Apply formatting
        self._apply_shape_formatting(shape, elem)

    def _apply_shape_formatting(self, shape, elem):
        """Apply fill and line formatting to a shape."""
        fill_color = elem.get("fill")
        if fill_color == "none":
            shape.fill.background()
        elif fill_color:
            shape.fill.solid()
            hex_color = fill_color.lstrip("#")
            shape.fill.fore_color.rgb = RGBColor(
                int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
            )
            opacity = elem.get("opacity")
            if opacity is not None and 0 <= opacity < 1:
                self._set_fill_opacity(shape, opacity)
        
        line_color = elem.get("line")
        line_width = elem.get("lineWidth", 1)
        line_gradient = elem.get("lineGradient")
        
        if line_gradient:
            from lxml import etree
            from pptx.oxml.ns import qn
            sp_pr = shape._element.spPr
            ln = sp_pr.find(qn('a:ln'))
            if ln is None:
                ln = etree.SubElement(sp_pr, qn('a:ln'))
            ln.set('w', str(int(line_width * 12700)))
            ln.set('cap', 'flat')
            # Remove existing fill
            for child in list(ln):
                if child.tag.endswith('}solidFill') or child.tag.endswith('}noFill') or child.tag.endswith('}gradFill'):
                    ln.remove(child)
            grad = etree.SubElement(ln, qn('a:gradFill'))
            gs_lst = etree.SubElement(grad, qn('a:gsLst'))
            for stop in line_gradient.get("stops", []):
                gs = etree.SubElement(gs_lst, qn('a:gs'))
                gs.set('pos', str(int(stop.get("position", 0) * 100000)))
                srgb = etree.SubElement(gs, qn('a:srgbClr'))
                srgb.set('val', stop.get("color", "#FFFFFF").lstrip("#"))
                opacity = stop.get("opacity")
                if opacity is not None and opacity < 1:
                    alpha = etree.SubElement(srgb, qn('a:alpha'))
                    alpha.set('val', str(int(opacity * 100000)))
            lin = etree.SubElement(grad, qn('a:lin'))
            lin.set('ang', str(int(line_gradient.get("angle", 0) * 60000)))
            lin.set('scaled', '1')
        elif line_color and line_color != "none":
            shape.line.fill.solid()
            hex_color = line_color.lstrip("#")
            shape.line.color.rgb = RGBColor(
                int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
            )
            shape.line.width = Pt(line_width)
            # Set line cap to flat for freeform shapes
            try:
                from pptx.oxml.ns import qn
                sp_pr = shape._element.spPr
                ln = sp_pr.find(qn('a:ln'))
                if ln is not None:
                    ln.set('cap', 'flat')
            except:
                pass
        elif line_color == "none":
            shape.line.fill.background()
    
    def _px_to_emu(self, px):
        """Convert pixels (1920x1080 basis) to EMU."""
        return Emu(int(px * self.EMU_PER_PX))
    
    def _set_fill_opacity(self, shape, opacity):
        """Set fill opacity using low-level XML manipulation.
        
        Args:
            shape: Shape object
            opacity: 0.0 (fully transparent) to 1.0 (fully opaque)
        """
        from lxml import etree
        # alpha value in OOXML is percentage * 1000 (e.g., 50% = 50000)
        alpha_val = int(opacity * 100000)
        solidFill = shape.fill._xPr.solidFill
        if solidFill is not None:
            srgbClr = solidFill.srgbClr
            if srgbClr is not None:
                # Remove existing alpha if present
                nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                for existing in srgbClr.findall('a:alpha', nsmap):
                    srgbClr.remove(existing)
                # Add alpha element
                alpha_elem = etree.SubElement(srgbClr, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
                alpha_elem.set('val', str(alpha_val))
    
    def _add_table(self, slide, elem):
        from pptx.enum.text import PP_ALIGN
        from lxml import etree

        headers = elem.get("headers", [])
        rows = elem.get("rows", [])
        cols = len(headers) if headers else (len(rows[0]) if rows else 0)
        row_count = len(rows) + (1 if headers else 0)

        if cols == 0 or row_count == 0:
            return

        x = self._px_to_emu(elem.get("x", 77))
        y = self._px_to_emu(elem.get("y", 270))
        width = self._px_to_emu(elem.get("width", 1766))
        height = self._px_to_emu(elem.get("height")) if elem.get("height") else Emu(row_count * 400000)

        tbl_shape = slide.shapes.add_table(row_count, cols, x, y, width, height)
        table = tbl_shape.table
        colors = self.colors
        nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

        # Column widths
        col_widths = elem.get("colWidths")
        if col_widths:
            for i, w in enumerate(col_widths):
                if i < len(table.columns):
                    table.columns[i].width = self._px_to_emu(w)

        # Row heights
        row_heights = elem.get("rowHeights")
        if row_heights:
            for i, h in enumerate(row_heights):
                if i < len(table.rows):
                    table.rows[i].height = self._px_to_emu(h)

        # Table style properties
        tbl_pr = table._tbl.find('a:tblPr', nsmap)
        if tbl_pr is not None:
            for attr in ['firstRow', 'lastRow', 'firstCol', 'lastCol', 'bandRow', 'bandCol']:
                if elem.get(attr):
                    tbl_pr.set(attr, '1')
                elif attr in tbl_pr.attrib:
                    del tbl_pr.attrib[attr]
            style_id = elem.get("tableStyleId")
            if style_id:
                existing = tbl_pr.find('a:tableStyleId', nsmap)
                if existing is not None:
                    existing.text = style_id
                else:
                    el = etree.SubElement(tbl_pr, f'{{{nsmap["a"]}}}tableStyleId')
                    el.text = style_id

        # Merge cells first (scan all rows for gridSpan/rowSpan)
        all_rows = [headers] + rows if headers else rows
        for ri, row_data in enumerate(all_rows):
            for ci, cell_val in enumerate(row_data):
                if not isinstance(cell_val, dict):
                    continue
                gs = cell_val.get("gridSpan", 1)
                rs = cell_val.get("rowSpan", 1)
                if gs > 1 or rs > 1:
                    try:
                        table.cell(ri, ci).merge(table.cell(ri + rs - 1, ci + gs - 1))
                    except Exception:
                        pass

        # Helper to apply cell properties
        def apply_cell(cell, cell_val, fallback_fill=None, fallback_text_color=None):
            text = cell_val if isinstance(cell_val, str) else cell_val.get("text", "")
            is_dict = isinstance(cell_val, dict)

            # Skip merged-away cells
            if is_dict and cell_val.get("merged"):
                return

            # Fill
            if is_dict and "fill" in cell_val:
                f = cell_val["fill"]
                if f and f != "none":
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor.from_string(f.lstrip('#'))
            elif fallback_fill:
                cell.fill.solid()
                cell.fill.fore_color.rgb = fallback_fill

            # Vertical alignment
            if is_dict and cell_val.get("anchor"):
                anchor_map = {'t': 0, 'ctr': 1, 'b': 2}  # TOP/MIDDLE/BOTTOM
                tc_pr = cell._tc.find(f'{{{nsmap["a"]}}}tcPr')
                if tc_pr is None:
                    tc_pr = etree.SubElement(cell._tc, f'{{{nsmap["a"]}}}tcPr')
                tc_pr.set('anchor', cell_val["anchor"])

            # Margins
            if is_dict and cell_val.get("margins"):
                tc_pr = cell._tc.find(f'{{{nsmap["a"]}}}tcPr')
                if tc_pr is None:
                    tc_pr = etree.SubElement(cell._tc, f'{{{nsmap["a"]}}}tcPr')
                m = cell_val["margins"]
                for side, attr in [('left', 'marL'), ('right', 'marR'), ('top', 'marT'), ('bottom', 'marB')]:
                    if side in m:
                        tc_pr.set(attr, str(self._px_to_emu(m[side])))

            # Borders
            if is_dict and cell_val.get("borders"):
                tc_pr = cell._tc.find(f'{{{nsmap["a"]}}}tcPr')
                if tc_pr is None:
                    tc_pr = etree.SubElement(cell._tc, f'{{{nsmap["a"]}}}tcPr')
                tag_map = {"left": "lnL", "right": "lnR", "top": "lnT", "bottom": "lnB"}
                for side, bdr in cell_val["borders"].items():
                    tag = tag_map.get(side)
                    if not tag:
                        continue
                    # Remove existing
                    existing = tc_pr.find(f'{{{nsmap["a"]}}}{tag}')
                    if existing is not None:
                        tc_pr.remove(existing)
                    ln = etree.SubElement(tc_pr, f'{{{nsmap["a"]}}}{tag}')
                    if bdr.get("width"):
                        ln.set('w', str(int(bdr["width"] * 12700)))
                    if bdr.get("fill") == "none":
                        etree.SubElement(ln, f'{{{nsmap["a"]}}}noFill')
                    elif bdr.get("color"):
                        sf = etree.SubElement(ln, f'{{{nsmap["a"]}}}solidFill')
                        srgb = etree.SubElement(sf, f'{{{nsmap["a"]}}}srgbClr')
                        srgb.set('val', bdr["color"].lstrip('#'))

            # Text
            para = cell.text_frame.paragraphs[0]
            text_color = fallback_text_color
            if is_dict:
                if cell_val.get("fontColor"):
                    text_color = RGBColor.from_string(cell_val["fontColor"].lstrip('#'))
                # Alignment
                align = cell_val.get("align")
                if align:
                    para.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(align)
                # Font size
                font_size = cell_val.get("fontSize")
            else:
                font_size = None

            self._apply_styled_text(para, str(text), default_color=text_color)

            # Apply bold/italic/fontSize on all runs
            if is_dict:
                for run in para.runs:
                    if cell_val.get("bold"):
                        run.font.bold = True
                    if cell_val.get("italic"):
                        run.font.italic = True
                    if font_size:
                        run.font.size = Pt(font_size)

        # Fill headers
        has_table_style = bool(elem.get("tableStyleId"))
        if headers:
            for ci, hdr in enumerate(headers):
                apply_cell(table.cell(0, ci), hdr,
                           fallback_fill=None if has_table_style else colors["header"],
                           fallback_text_color=None if has_table_style else colors["text_inverse"])

        # Fill data rows
        start_row = 1 if headers else 0
        for ri, row in enumerate(rows):
            bg = colors["background_alt"] if ri % 2 == 1 else colors["background"]
            for ci, val in enumerate(row):
                apply_cell(table.cell(start_row + ri, ci), val,
                           fallback_fill=None if has_table_style else bg,
                           fallback_text_color=None if has_table_style else colors["text"])
    
    def _add_textbox(self, slide, elem):
        from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
        
        align = elem.get("align", "left")
        x_pct = elem.get("x", 3)
        y_pct = elem.get("y", 25)
        width_pct = elem.get("width", 94)
        auto_width = elem.get("autoWidth", False)
        
        # Convert percentage to EMU (x is always left edge)
        x_emu = self._px_to_emu(x_pct)
        y_emu = self._px_to_emu(y_pct)
        width_emu = self._px_to_emu(width_pct)
        height_px = elem.get("height")
        height_emu = self._px_to_emu(height_px) if height_px else self._px_to_emu(10)
        
        textbox = slide.shapes.add_textbox(x_emu, y_emu, width_emu, height_emu)
        tf = textbox.text_frame
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE if height_px else MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        tf.word_wrap = not auto_width
        
        # Apply rotation
        rotation = elem.get("rotation", 0)
        if rotation != 0:
            textbox.rotation = rotation
        
        # Apply flip via XML
        flip_h = elem.get("flipH", False)
        flip_v = elem.get("flipV", False)
        
        if flip_h or flip_v:
            try:
                from lxml import etree
                xfrm = textbox._element.spPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm')
                if xfrm is not None:
                    if flip_h:
                        xfrm.set('flipH', '1')
                    if flip_v:
                        xfrm.set('flipV', '1')
            except:
                pass
        
        # Apply fill color
        fill_color = elem.get("fill")
        if fill_color == "none" or fill_color is None:
            textbox.fill.background()
        elif fill_color:
            textbox.fill.solid()
            hex_color = fill_color.lstrip("#")
            textbox.fill.fore_color.rgb = RGBColor(
                int(hex_color[0:2], 16),
                int(hex_color[2:4], 16),
                int(hex_color[4:6], 16)
            )
        
        # Set paragraph alignment
        align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
        
        # Check if paragraphs array exists (multiple paragraphs without bullets)
        paragraphs = elem.get("paragraphs")
        
        if paragraphs:
            # Add multiple paragraphs (with or without bullets/numbering)
            default_font_size = elem.get("fontSize")
            for i, para_item in enumerate(paragraphs):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                # Check if it's a dict with bullet/numbering info or just a string
                if isinstance(para_item, dict):
                    para_text = para_item.get("text", "")
                    has_bullet = para_item.get("bullet", False)
                    numbering = para_item.get("numbering")
                else:
                    para_text = para_item
                    has_bullet = False
                    numbering = None
                
                # Apply bullet or numbering
                if numbering:
                    p.level = 0
                    self._set_numbering(p, numbering)
                elif has_bullet:
                    p.level = 0
                    self._set_bullet(p)
                
                # Apply text
                if para_text:
                    self._apply_styled_text(p, para_text, default_font_size=default_font_size)
                
                # Apply space after
                if isinstance(para_item, dict) and para_item.get("spaceAfter") is not None:
                    p.space_after = Pt(para_item["spaceAfter"] / 100)
                
                p.alignment = align_map.get(align, PP_ALIGN.LEFT)
        else:
            # Single text — expand styled tags spanning newlines, then split
            text = elem.get("text", "")
            text = _expand_styled_newlines(text)
            default_font_size = elem.get("fontSize")
            lines = text.split("\n")
            for i, line in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                self._apply_styled_text(p, line, default_font_size=default_font_size)
                p.alignment = align_map.get(align, PP_ALIGN.LEFT)
        
        # Apply line (border) to textbox
        line_color = elem.get("line")
        line_width = elem.get("lineWidth", 1)
        
        if line_color == "none":
            textbox.line.fill.background()
        elif line_color:
            textbox.line.fill.solid()
            hex_color = line_color.lstrip("#")
            textbox.line.color.rgb = RGBColor(
                int(hex_color[0:2], 16),
                int(hex_color[2:4], 16),
                int(hex_color[4:6], 16)
            )
            textbox.line.width = Pt(line_width)
        
        # Apply text frame margins
        margin_left = elem.get("marginLeft")
        if margin_left is not None:
            tf.margin_left = margin_left
        margin_top = elem.get("marginTop")
        if margin_top is not None:
            tf.margin_top = margin_top
        margin_right = elem.get("marginRight")
        if margin_right is not None:
            tf.margin_right = margin_right
        margin_bottom = elem.get("marginBottom")
        if margin_bottom is not None:
            tf.margin_bottom = margin_bottom
        
        # Apply vertical anchor
        vertical_anchor = elem.get("verticalAnchor")
        if vertical_anchor is not None:
            tf.vertical_anchor = vertical_anchor
        
        # Apply text gradient
        text_gradient = elem.get("textGradient")
        if text_gradient:
            self._apply_text_gradient(textbox, text_gradient)
    
    def _add_line(self, slide, elem):
        """Add line/connector to slide."""
        x_pct = elem.get("x", 10)
        y_pct = elem.get("y", 10)
        width_pct = elem.get("width", 20)
        height_pct = elem.get("height", 0.1)
        
        x_emu = self._px_to_emu(x_pct)
        y_emu = self._px_to_emu(y_pct)
        width_emu = self._px_to_emu(width_pct)
        height_emu = self._px_to_emu(height_pct)
        
        # Determine connector type
        connector_type_str = elem.get("connectorType", "straight")
        connector_type_map = {
            "straight": MSO_CONNECTOR.STRAIGHT,
            "elbow": MSO_CONNECTOR.ELBOW,
            "curved": MSO_CONNECTOR.CURVE,
        }
        connector_type = connector_type_map.get(connector_type_str, MSO_CONNECTOR.STRAIGHT)
        
        # Add connector (line)
        connector = slide.shapes.add_connector(
            connector_type,
            x_emu, y_emu,
            x_emu + width_emu, y_emu + height_emu
        )
        
        # Apply rotation
        rotation = elem.get("rotation", 0)
        if rotation != 0:
            connector.rotation = rotation
        
        # Apply flip via XML
        flip_h = elem.get("flipH", False)
        flip_v = elem.get("flipV", False)
        
        # AWS recommended defaults (theme-aware)
        default_color = "#8FA7C4" if self.theme == "dark" else "#000000"
        default_line_width = 1.25
        
        if flip_h or flip_v:
            try:
                from lxml import etree
                xfrm = connector._element.spPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm')
                if xfrm is not None:
                    if flip_h:
                        xfrm.set('flipH', '1')
                    if flip_v:
                        xfrm.set('flipV', '1')
            except:
                pass
        
        # Apply exact preset geometry if available
        preset = elem.get("preset")
        if preset:
            try:
                from lxml import etree
                sp_pr = connector._element.spPr
                prst_geom = sp_pr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
                if prst_geom is not None:
                    prst_geom.set('prst', preset)
            except:
                pass
        
        # Apply arrow heads via XML
        head_end = elem.get("headEnd")
        tail_end = elem.get("tailEnd")
        
        if head_end or tail_end:
            from lxml import etree
            ln = connector._element.spPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}ln')
            if ln is None:
                ln = etree.SubElement(connector._element.spPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}ln')
            
            if head_end:
                head_elem = etree.SubElement(ln, '{http://schemas.openxmlformats.org/drawingml/2006/main}headEnd')
                head_elem.set('type', head_end)
            
            if tail_end:
                tail_elem = etree.SubElement(ln, '{http://schemas.openxmlformats.org/drawingml/2006/main}tailEnd')
                tail_elem.set('type', tail_end)
        
        # Apply line color or gradient
        line_gradient = elem.get("lineGradient")
        color = elem.get("color", default_color)
        line_width = elem.get("lineWidth", default_line_width)
        
        if line_gradient:
            # Apply line gradient
            try:
                connector.line.fill.gradient()
                stops = line_gradient.get("stops", [])
                
                # Update gradient stops - need to handle more than 2 stops via XML
                gradient_stops = connector.line.fill.gradient_stops
                
                if len(stops) > 2 or any(s.get("opacity") is not None for s in stops):
                    # Add additional stops via XML (to support 3+ stops and opacity)
                    from lxml import etree
                    gsLst = connector.line._ln.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}gsLst')
                    if gsLst is not None:
                        # Clear existing stops
                        for gs in list(gsLst):
                            gsLst.remove(gs)
                        # Add all stops
                        for stop_info in stops:
                            pos = int(stop_info.get("position", 0) * 100000)
                            color_hex = stop_info.get("color", "#FFFFFF").lstrip("#")
                            gs = etree.SubElement(gsLst, '{http://schemas.openxmlformats.org/drawingml/2006/main}gs')
                            gs.set('pos', str(pos))
                            srgbClr = etree.SubElement(gs, '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
                            srgbClr.set('val', color_hex)
                            # Add opacity if specified
                            stop_opacity = stop_info.get("opacity")
                            if stop_opacity is not None and stop_opacity < 1:
                                alpha = etree.SubElement(srgbClr, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
                                alpha.set('val', str(int(stop_opacity * 100000)))
                else:
                    for i, stop_info in enumerate(stops):
                        if i < len(gradient_stops):
                            pos = stop_info.get("position", 0)
                            color_hex = stop_info.get("color", "#FFFFFF")
                            hex_color = color_hex.lstrip("#")
                            stop_color = RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))
                            
                            gradient_stops[i].color.rgb = stop_color
                            gradient_stops[i].position = pos
                
                # Try to set angle via XML
                angle = line_gradient.get("angle", 0)
                try:
                    from lxml import etree
                    lin = connector.line._ln.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}lin')
                    if lin is not None:
                        # Convert angle to EMUs (60000 per degree)
                        lin.set('ang', str(int(angle * 60000)))
                    else:
                        # Create lin element if not exists
                        gradFill = connector.line._ln.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}gradFill')
                        if gradFill is not None:
                            lin = etree.SubElement(gradFill, '{http://schemas.openxmlformats.org/drawingml/2006/main}lin')
                            lin.set('ang', str(int(angle * 60000)))
                except:
                    pass
                
                connector.line.width = Pt(line_width)
            except:
                # Fallback to solid color
                if stops:
                    color_hex = stops[0].get("color", "#FFFFFF")
                    hex_color = color_hex.lstrip("#")
                    connector.line.fill.solid()
                    connector.line.color.rgb = RGBColor(
                        int(hex_color[0:2], 16),
                        int(hex_color[2:4], 16),
                        int(hex_color[4:6], 16)
                    )
                    connector.line.width = Pt(line_width)
        elif color:
            connector.line.fill.solid()
            hex_color = color.lstrip("#")
            connector.line.color.rgb = RGBColor(
                int(hex_color[0:2], 16),
                int(hex_color[2:4], 16),
                int(hex_color[4:6], 16)
            )
            connector.line.width = Pt(line_width)
        
        # Apply dash style
        dash_style = elem.get("dashStyle")
        if dash_style:
            from pptx.enum.dml import MSO_LINE_DASH_STYLE
            dash_map = {
                "solid": MSO_LINE_DASH_STYLE.SOLID,
                "dash": MSO_LINE_DASH_STYLE.DASH,
                "dot": MSO_LINE_DASH_STYLE.ROUND_DOT,
                "dash_dot": MSO_LINE_DASH_STYLE.DASH_DOT,
                "dash_dot_dot": MSO_LINE_DASH_STYLE.DASH_DOT_DOT,
                "long_dash": MSO_LINE_DASH_STYLE.LONG_DASH,
                "long_dash_dot": MSO_LINE_DASH_STYLE.LONG_DASH_DOT,
                "square_dot": MSO_LINE_DASH_STYLE.SQUARE_DOT,
            }
            if dash_style in dash_map:
                connector.line.dash_style = dash_map[dash_style]
    
    def _add_shape(self, slide, elem):
        """Add shape to slide.
        
        Supported shapes:
        - rectangle, rounded_rectangle
        - oval, circle
        - arrow_right, arrow_left, arrow_up, arrow_down
        - triangle, diamond, pentagon, hexagon
        
        Args:
            shape: Shape type (required)
            x, y: Position in percentage (required)
            width, height: Size in percentage (required)
            fill: Fill color (#RRGGBB or "none")
            line: Line color (#RRGGBB or "none")
            lineWidth: Line width in pt (default: 1)
            text: Text inside shape (optional)
            fontSize: Font size for text (default: 14)
        """
        shape_type = elem.get("shape")
        if not shape_type:
            return
        
        # Map shape names to MSO_SHAPE constants
        shape_map = {
            "rectangle": MSO_SHAPE.RECTANGLE,
            "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
            "oval": MSO_SHAPE.OVAL,
            "circle": MSO_SHAPE.OVAL,  # Circle is oval with equal width/height
            "arrow_right": MSO_SHAPE.RIGHT_ARROW,
            "arrow_left": MSO_SHAPE.LEFT_ARROW,
            "arrow_up": MSO_SHAPE.UP_ARROW,
            "arrow_down": MSO_SHAPE.DOWN_ARROW,
            "arrow_circular": MSO_SHAPE.CIRCULAR_ARROW,
            "arrow_left_right": MSO_SHAPE.LEFT_RIGHT_ARROW,
            "arrow_up_down": MSO_SHAPE.UP_DOWN_ARROW,
            "arrow_curved_right": MSO_SHAPE.CURVED_RIGHT_ARROW,
            "arrow_curved_left": MSO_SHAPE.CURVED_LEFT_ARROW,
            "arrow_curved_up": MSO_SHAPE.CURVED_UP_ARROW,
            "arrow_curved_down": MSO_SHAPE.CURVED_DOWN_ARROW,
            "arrow_circular_left": MSO_SHAPE.LEFT_CIRCULAR_ARROW,
            "arrow_circular_left_right": MSO_SHAPE.LEFT_RIGHT_CIRCULAR_ARROW,
            "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
            "diamond": MSO_SHAPE.DIAMOND,
            "pentagon": MSO_SHAPE.PENTAGON,
            "hexagon": MSO_SHAPE.HEXAGON,
            "cross": MSO_SHAPE.CROSS,
            "trapezoid": MSO_SHAPE.TRAPEZOID,
            "parallelogram": MSO_SHAPE.PARALLELOGRAM,
            "chevron": MSO_SHAPE.CHEVRON,
            "donut": MSO_SHAPE.DONUT,
            "arc": MSO_SHAPE.ARC,
            "block_arc": MSO_SHAPE.BLOCK_ARC,
            "chord": MSO_SHAPE.CHORD,
            "pie": MSO_SHAPE.PIE,
            "pie_wedge": MSO_SHAPE.PIE_WEDGE,
            "cloud": MSO_SHAPE.CLOUD,
            "lightning_bolt": MSO_SHAPE.LIGHTNING_BOLT,
            "star_5_point": MSO_SHAPE.STAR_5_POINT,
            "no_symbol": MSO_SHAPE.NO_SYMBOL,
            "callout_rectangle": MSO_SHAPE.RECTANGULAR_CALLOUT,
            "callout_rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGULAR_CALLOUT,
            "callout_oval": MSO_SHAPE.OVAL_CALLOUT,
            "flowchart_process": MSO_SHAPE.FLOWCHART_PROCESS,
            "flowchart_decision": MSO_SHAPE.FLOWCHART_DECISION,
            "flowchart_terminator": MSO_SHAPE.FLOWCHART_TERMINATOR,
            "left_brace": MSO_SHAPE.LEFT_BRACE,
            "right_brace": MSO_SHAPE.RIGHT_BRACE,
            "left_bracket": MSO_SHAPE.LEFT_BRACKET,
            "right_bracket": MSO_SHAPE.RIGHT_BRACKET,
        }
        
        mso_shape = shape_map.get(shape_type)
        if not mso_shape:
            print(f"Warning: Unknown shape: {shape_type}", file=sys.stderr)
            return
        
        x_pct = elem.get("x", 10)
        y_pct = elem.get("y", 10)
        width_pct = elem.get("width", 20)
        height_pct = elem.get("height", 10)
        
        x_emu = self._px_to_emu(x_pct)
        y_emu = self._px_to_emu(y_pct)
        width_emu = self._px_to_emu(width_pct)
        height_emu = self._px_to_emu(height_pct)
        
        # For circle, use minimum of width/height
        if shape_type == "circle":
            size = min(width_emu, height_emu)
            width_emu = height_emu = size
        
        shape = slide.shapes.add_shape(mso_shape, x_emu, y_emu, width_emu, height_emu)
        
        # Apply rotation
        rotation = elem.get("rotation", 0)
        if rotation != 0:
            shape.rotation = rotation
        
        # Apply adjustments (shape control points)
        adjustments = elem.get("adjustments")
        if adjustments and hasattr(shape, 'adjustments'):
            try:
                for i, adj_val in enumerate(adjustments):
                    if i < len(shape.adjustments):
                        shape.adjustments[i] = adj_val
            except:
                pass
        
        # Apply fill color or gradient
        gradient = elem.get("gradient")
        fill_color = elem.get("fill")
        
        if gradient:
            # Apply gradient
            shape.fill.gradient()
            stops = gradient.get("stops", [])
            angle = gradient.get("angle", 0)
            
            # Set gradient angle
            shape.fill.gradient_angle = angle
            
            # Update gradient stops via XML (to support 3+ stops and opacity)
            try:
                from lxml import etree
                from pptx.oxml.ns import qn
                
                sp_pr = shape._element.spPr
                grad_fill = sp_pr.find(qn('a:gradFill'))
                if grad_fill is not None:
                    gs_lst = grad_fill.find(qn('a:gsLst'))
                    if gs_lst is not None:
                        # Remove all existing stops
                        for gs in list(gs_lst):
                            gs_lst.remove(gs)
                        
                        # Add all stops from JSON
                        for stop_info in stops:
                            pos = int(stop_info.get("position", 0) * 100000)
                            color_hex = stop_info.get("color", "#FFFFFF").lstrip("#")
                            stop_opacity = stop_info.get("opacity")
                            
                            new_gs = etree.Element(qn('a:gs'))
                            new_gs.set('pos', str(pos))
                            
                            srgb_clr = etree.SubElement(new_gs, qn('a:srgbClr'))
                            srgb_clr.set('val', color_hex.upper())
                            
                            # Add opacity if needed
                            if stop_opacity is not None and stop_opacity < 1:
                                alpha = etree.SubElement(srgb_clr, qn('a:alpha'))
                                alpha.set('val', str(int(stop_opacity * 100000)))
                            
                            gs_lst.append(new_gs)
            except Exception as e:
                print(f"Warning: Failed to update gradient stops: {e}", file=sys.stderr)
        elif fill_color == "none" or fill_color is None:
            # Transparent fill
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
            shape.fill.background()  # Set to transparent
        elif fill_color:
            shape.fill.solid()
            hex_color = fill_color.lstrip("#")
            shape.fill.fore_color.rgb = RGBColor(
                int(hex_color[0:2], 16),
                int(hex_color[2:4], 16),
                int(hex_color[4:6], 16)
            )
            # Apply fill opacity (0.0 = fully transparent, 1.0 = fully opaque)
            opacity = elem.get("opacity")
            if opacity is not None and 0 <= opacity < 1:
                self._set_fill_opacity(shape, opacity)
        
        # Apply line color or gradient
        line_gradient = elem.get("lineGradient")
        line_color = elem.get("line")
        line_width = elem.get("lineWidth", 1)
        
        if line_gradient:
            # Apply line gradient
            try:
                shape.line.fill.gradient()
                stops = line_gradient.get("stops", [])
                
                # Update gradient stops - need to handle more than 2 stops via XML
                gradient_stops = shape.line.fill.gradient_stops
                
                if len(stops) > 2:
                    # Add additional stops via XML
                    from lxml import etree
                    gsLst = shape.line._ln.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}gsLst')
                    if gsLst is not None:
                        # Clear existing stops
                        for gs in list(gsLst):
                            gsLst.remove(gs)
                        # Add all stops
                        for stop_info in stops:
                            pos = int(stop_info.get("position", 0) * 100000)
                            color_hex = stop_info.get("color", "#FFFFFF").lstrip("#")
                            gs = etree.SubElement(gsLst, '{http://schemas.openxmlformats.org/drawingml/2006/main}gs')
                            gs.set('pos', str(pos))
                            srgbClr = etree.SubElement(gs, '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
                            srgbClr.set('val', color_hex)
                else:
                    for i, stop_info in enumerate(stops):
                        if i < len(gradient_stops):
                            pos = stop_info.get("position", 0)
                            color_hex = stop_info.get("color", "#FFFFFF")
                            hex_color = color_hex.lstrip("#")
                            color = RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))
                            
                            gradient_stops[i].color.rgb = color
                            gradient_stops[i].position = pos
                
                # Set angle via XML
                angle = line_gradient.get("angle", 0)
                try:
                    from lxml import etree
                    gradFill = shape.line._ln.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}gradFill')
                    if gradFill is not None:
                        lin = gradFill.find('{http://schemas.openxmlformats.org/drawingml/2006/main}lin')
                        if lin is None:
                            lin = etree.SubElement(gradFill, '{http://schemas.openxmlformats.org/drawingml/2006/main}lin')
                        lin.set('ang', str(int(angle * 60000)))
                        lin.set('scaled', '1')
                except:
                    pass
                
                shape.line.width = Pt(line_width)
            except Exception as e:
                # Fallback to solid color (use first stop color)
                print(f"Warning: Line gradient failed, using solid color: {e}", file=sys.stderr)
                if stops:
                    color_hex = stops[0].get("color", "#8B8BF8")
                    hex_color = color_hex.lstrip("#")
                    shape.line.fill.solid()
                    shape.line.color.rgb = RGBColor(
                        int(hex_color[0:2], 16),
                        int(hex_color[2:4], 16),
                        int(hex_color[4:6], 16)
                    )
                    shape.line.width = Pt(line_width)
        elif line_color == "none":
            shape.line.fill.background()
        elif line_color:
            shape.line.fill.solid()
            hex_color = line_color.lstrip("#")
            shape.line.color.rgb = RGBColor(
                int(hex_color[0:2], 16),
                int(hex_color[2:4], 16),
                int(hex_color[4:6], 16)
            )
            shape.line.width = Pt(line_width)
        
        # Apply dash style to shape line
        dash_style = elem.get("dashStyle")
        if dash_style:
            from pptx.enum.dml import MSO_LINE_DASH_STYLE
            dash_map = {
                "solid": MSO_LINE_DASH_STYLE.SOLID,
                "dash": MSO_LINE_DASH_STYLE.DASH,
                "dot": MSO_LINE_DASH_STYLE.ROUND_DOT,
                "dash_dot": MSO_LINE_DASH_STYLE.DASH_DOT,
                "dash_dot_dot": MSO_LINE_DASH_STYLE.DASH_DOT_DOT,
                "long_dash": MSO_LINE_DASH_STYLE.LONG_DASH,
                "long_dash_dot": MSO_LINE_DASH_STYLE.LONG_DASH_DOT,
                "square_dot": MSO_LINE_DASH_STYLE.SQUARE_DOT,
            }
            if dash_style in dash_map:
                shape.line.dash_style = dash_map[dash_style]
        
        # Add text if specified
        text = elem.get("text")
        items = elem.get("items")
        
        if (text or items) and shape.has_text_frame:
            from pptx.enum.text import MSO_AUTO_SIZE
            tf = shape.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            tf.clear()
            
            # Apply vertical anchor (default: middle=3)
            from pptx.enum.text import MSO_ANCHOR
            vertical_anchor = elem.get("verticalAnchor")
            if vertical_anchor is not None:
                tf.vertical_anchor = vertical_anchor
            else:
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            font_size = elem.get("fontSize", 14)
            
            if items:
                # Add bullet points
                for i, item in enumerate(items):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.level = 0
                    self._set_bullet(p)
                    item_text = item.get("text", item) if isinstance(item, dict) else item
                    self._apply_styled_text(p, item_text, default_font_size=font_size)
                    
                    if isinstance(item, dict) and item.get("spaceAfter") is not None:
                        p.space_after = Pt(item["spaceAfter"] / 100)
                    
                    # Apply text alignment
                    text_align = elem.get("textAlign", "center")
                    if text_align == "center":
                        p.alignment = 2
                    elif text_align == "right":
                        p.alignment = 3
                    elif text_align == "left":
                        p.alignment = 1
                    else:
                        p.alignment = 2  # Default center
            else:
                # Single text
                self._apply_styled_text(tf.paragraphs[0], text, default_font_size=font_size)
                
                # Apply text alignment
                text_align = elem.get("textAlign", "center")
                if text_align == "center":
                    tf.paragraphs[0].alignment = 2  # CENTER
                elif text_align == "right":
                    tf.paragraphs[0].alignment = 3  # RIGHT
                elif text_align == "left":
                    tf.paragraphs[0].alignment = 1  # LEFT
                else:
                    tf.paragraphs[0].alignment = 2  # Default center for shapes
        
        # Add hyperlink to entire shape if specified
        link = elem.get("link")
        if link:
            shape.click_action.hyperlink.address = link
    
    def _add_image(self, slide, elem):
        """Add image element to slide.
        
        src: icons:NAME or file path (supports ~)
        """
        from pptx.enum.text import PP_ALIGN
        
        src = elem.get("src") or elem.get("path", "")
        x_pct = elem.get("x", 0)
        y_pct = elem.get("y", 0)
        width_pct = elem.get("width")
        height_pct = elem.get("height")
        label = elem.get("label")
        label_pos = elem.get("labelPosition", "bottom")
        label_size = elem.get("labelSize", 11)
        link = elem.get("link")
        rotation = elem.get("rotation", 0)
        icon_color = elem.get("iconColor")
        
        if not src:
            return
        
        # Resolve image path
        if src.startswith("icons:"):
            img_path = resolve_image_path(src, self.theme)
        else:
            img_path = Path(src).expanduser()
            if not img_path.is_absolute():
                img_path = Path.cwd() / src
            if not img_path.exists():
                print(f"Warning: Image not found: {img_path}", file=sys.stderr)
                return
        
        is_svg = img_path.suffix.lower() == ".svg"
        
        # Prepare SVG bytes (with optional recolor)
        svg_bytes = None
        if is_svg:
            svg_bytes = img_path.read_bytes()
            if icon_color:
                recolored = _recolor_svg(svg_bytes, icon_color)
                if recolored:
                    svg_bytes = recolored
        elif icon_color:
            print(f"Warning: iconColor ignored (not SVG): {img_path.name}", file=sys.stderr)
        
        # Calculate dimensions
        x = self._px_to_emu(x_pct)
        y = self._px_to_emu(y_pct)
        
        if width_pct:
            width = self._px_to_emu(width_pct)
            if height_pct:
                height = self._px_to_emu(height_pct)
            else:
                # Maintain original aspect ratio
                if is_svg:
                    img_w, img_h = get_svg_dimensions(img_path)
                else:
                    from PIL import Image
                    try:
                        with Image.open(img_path) as img:
                            img_w, img_h = img.size
                    except:
                        img_w, img_h = 1, 1
                height = int(width * img_h / img_w) if img_w > 0 else width
        
        if is_svg:
            if not width_pct:
                img_w, img_h = get_svg_dimensions(img_path)
                width = self._px_to_emu(img_w)
                height = self._px_to_emu(img_h)
            pic = add_svg_to_slide(slide, svg_bytes, x, y, width, height)
        else:
            if width_pct:
                pic = slide.shapes.add_picture(str(img_path), x, y, width=width, height=height)
            else:
                pic = slide.shapes.add_picture(str(img_path), x, y)
                width = pic.width
                height = pic.height
        
        # Apply rotation
        if rotation != 0:
            if is_svg:
                # Set rotation directly on XML xfrm element
                xfrm = pic.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm')
                if xfrm is not None:
                    xfrm.set('rot', str(int(rotation * 60000)))
            else:
                pic.rotation = rotation
        
        # Add hyperlink if specified
        if link:
            if is_svg:
                pass  # Hyperlinks on SVG not supported via direct XML (rare use case)
            else:
                pic.click_action.hyperlink.address = link
        
        # Add label if specified
        if label and label_pos != "none":
            if label_pos == "bottom":
                lbl_x = x
                lbl_y = y + height + Emu(50000)
                lbl_w = width
            elif label_pos == "right":
                lbl_x = x + width + Emu(100000)
                lbl_y = y + height // 3
                lbl_w = self._px_to_emu(15)
            else:
                return
            
            textbox = slide.shapes.add_textbox(lbl_x, lbl_y, lbl_w, Emu(300000))
            tf = textbox.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if label_pos == "bottom" else PP_ALIGN.LEFT
            self._apply_styled_text(p, label, default_font_size=label_size)
    
    def _apply_styled_text(self, paragraph, text, default_color=None, default_font_size=None):
        """Apply styled text to a paragraph.
        
        Supports:
        - {{bold:text}}
        - {{italic:text}}
        - {{#RRGGBB:text}}
        - {{NNpt:text}}
        - {{link:URL:text}} - hyperlink
        """
        segments = parse_styled_text(text)
        paragraph.clear()
        from lxml import etree
        from pptx.oxml.ns import qn
        
        for seg in segments:
            # Split text by character width for font assignment
            sub_runs = self._split_by_width(seg["text"])
            
            for sub_text, font_name in sub_runs:
                # Handle line breaks (\n and \u000b)
                lines = sub_text.replace('\u000b', '\n').split('\n')
                for li, line in enumerate(lines):
                    if li > 0:
                        etree.SubElement(paragraph._p, qn('a:br'))
                    run = paragraph.add_run()
                    run.text = line
                    run.font.name = font_name
                
                # Apply color (default or specified)
                if "color" in seg:
                    hex_color = seg["color"].lstrip("#")
                    color = RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))
                    run.font.color.rgb = color
                elif "link" not in seg:
                    color = default_color or self.colors["text"]
                    run.font.color.rgb = color
                
                # Apply font size
                font_size = seg.get("fontSize") or default_font_size
                if font_size:
                    run.font.size = Pt(font_size)
                
                # Apply bold/italic
                if seg.get("bold"):
                    run.font.bold = True
                if seg.get("italic"):
                    run.font.italic = True
                
                # Apply hyperlink (links use default blue color)
                if "link" in seg:
                    run.hyperlink.address = seg["link"]
                    run.font.underline = True

    def _apply_text_gradient(self, shape, gradient):
        """Apply gradient fill to text in a shape."""
        from lxml import etree
        from pptx.oxml.ns import qn
        
        stops = gradient.get("stops", [])
        angle = gradient.get("angle", 0)
        
        # Find all run elements and apply gradient
        tx_body = shape._element.find(qn('p:txBody'))
        if tx_body is None:
            return
        
        for r in tx_body.findall('.//' + qn('a:r')):
            rPr = r.find(qn('a:rPr'))
            if rPr is None:
                rPr = etree.Element(qn('a:rPr'))
                r.insert(0, rPr)
            
            # Remove existing solidFill and gradFill
            for sf in rPr.findall(qn('a:solidFill')):
                rPr.remove(sf)
            for gf in rPr.findall(qn('a:gradFill')):
                rPr.remove(gf)
            
            # Create gradFill element
            grad_fill = etree.Element(qn('a:gradFill'))
            gs_lst = etree.SubElement(grad_fill, qn('a:gsLst'))
            
            for stop_info in stops:
                pos = int(stop_info.get("position", 0) * 100000)
                color_hex = stop_info.get("color", "#FFFFFF").lstrip("#")
                stop_opacity = stop_info.get("opacity")
                
                gs = etree.SubElement(gs_lst, qn('a:gs'))
                gs.set('pos', str(pos))
                srgb_clr = etree.SubElement(gs, qn('a:srgbClr'))
                srgb_clr.set('val', color_hex.upper())
                
                # Add opacity if specified
                if stop_opacity is not None and stop_opacity < 1:
                    alpha = etree.SubElement(srgb_clr, qn('a:alpha'))
                    alpha.set('val', str(int(stop_opacity * 100000)))
            
            # Add linear gradient with angle
            lin = etree.SubElement(grad_fill, qn('a:lin'))
            lin.set('ang', str(int(angle * 60000)))
            lin.set('scaled', '0')
            
            # Insert gradFill at the beginning of rPr (before latin, ea, cs)
            rPr.insert(0, grad_fill)
    
    def _split_by_width(self, text):
        """Split text into runs by character width (fullwidth/halfwidth)."""
        if not text:
            return []
        
        runs = []
        current = []
        current_is_full = is_fullwidth(text[0])
        
        for char in text:
            char_is_full = is_fullwidth(char)
            if char_is_full == current_is_full:
                current.append(char)
            else:
                runs.append(("".join(current), FONT_FULLWIDTH if current_is_full else FONT_HALFWIDTH))
                current = [char]
                current_is_full = char_is_full
        
        if current:
            runs.append(("".join(current), FONT_FULLWIDTH if current_is_full else FONT_HALFWIDTH))
        
        return runs
    
    def save(self, output_path: Path):
        self.prs.save(str(output_path))


def resolve_override(slide_def, id_map, visited=None):
    """Resolve override chain and merge elements."""
    if visited is None:
        visited = set()
    
    override_id = slide_def.get("override")
    if not override_id:
        return slide_def
    
    # Circular reference check
    if override_id in visited:
        chain = " -> ".join(visited) + f" -> {override_id}"
        raise ValueError(f"Circular override detected: {chain}")
    
    if override_id not in id_map:
        raise ValueError(f"Override target '{override_id}' not found")
    
    visited.add(override_id)
    
    # Recursively resolve base (for chained overrides)
    base = resolve_override(id_map[override_id], id_map, visited)
    
    # Merge: base elements + override elements (override on top)
    base_elements = list(base.get("elements", []))
    override_elements = slide_def.get("elements", [])
    
    # Build result: copy override slide, replace elements with merged
    result = dict(slide_def)
    result["elements"] = base_elements + override_elements
    
    # Remove override key from result (already resolved)
    result.pop("override", None)
    
    return result


def validate_icons_in_json(data: dict, theme: str = "light") -> list:
    """Validate all icons in JSON and return list of missing icons."""
    missing = []
    slides = data.get("slides", [])
    
    def check_elements(elements):
        if not elements:
            return
        for elem in elements:
            if elem.get("type") == "image":
                src = elem.get("src") or elem.get("path", "")
                if src.startswith("icons:"):
                    icon_name = src[6:]
                    if not check_icon_exists(icon_name, theme):
                        missing.append(icon_name)
            elif elem.get("type") == "group":
                check_elements(elem.get("elements", []))
    
    for slide in slides:
        check_elements(slide.get("elements", []))
    
    return list(set(missing))  # Remove duplicates


def _is_wsl():
    return Path("/proc/version").exists() and "microsoft" in Path("/proc/version").read_text().lower()


def _refresh_autofit(pptx_path):
    """Refresh autofit by opening in PowerPoint and nudging a shape to trigger recalculation."""
    import subprocess
    if sys.platform == "darwin":
        script = '''
        on run argv
            set filePath to POSIX file (item 1 of argv) as alias
            tell application "Microsoft PowerPoint"
                activate
                open filePath
                delay 2
                set pres to active presentation
                if (count of shapes of slide 1 of pres) > 0 then
                    set sh to shape 1 of slide 1 of pres
                    set w to width of sh
                    set width of sh to w + 1
                    set width of sh to w
                end if
                delay 1
                save pres
                close pres
            end tell
        end run
        '''
        cmd = ["osascript", "-e", script, str(pptx_path)]
    elif _is_wsl():
        win_path = subprocess.run(["wslpath", "-w", str(pptx_path)], capture_output=True, text=True).stdout.strip()
        ps_cmd = (
            f"$app = New-Object -ComObject PowerPoint.Application; "
            f"$prs = $app.Presentations.Open('{win_path}'); "
            f"$s = $prs.Slides[1]; "
            f"if ($s.Shapes.Count -gt 0) {{ $sh = $s.Shapes[1]; $w = $sh.Width; $sh.Width = $w + 1; $sh.Width = $w }}; "
            f"$prs.Save(); $prs.Close(); $app.Quit()"
        )
        cmd = ["/mnt/c/Windows/System32/WindowsPowerShell/v1.0/powershell.exe", "-Command", ps_cmd]
    else:
        print("Warning: Autofit refresh skipped (unsupported platform)", file=sys.stderr)
        return
    try:
        result = subprocess.run(cmd, capture_output=True, timeout=30)
        if result.returncode == 0:
            print("Autofit refreshed via PowerPoint")
        else:
            encoding = "cp932" if _is_wsl() else "utf-8"
            stderr = result.stderr.decode(encoding, errors="replace").strip()
            print(f"Warning: Autofit refresh failed: {stderr}", file=sys.stderr)
    except (subprocess.TimeoutExpired, FileNotFoundError, UnicodeDecodeError):
        print("Warning: Autofit refresh skipped (PowerPoint not available)", file=sys.stderr)


def _unlock_height_constraints(pptx_path):
    """After autofit, bake scaled font sizes into runs and remove height constraints."""
    prs = Presentation(str(pptx_path))
    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    changed = False
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            body_pr = shape.text_frame._txBody.find('.//a:bodyPr', nsmap)
            if body_pr is None:
                continue
            norm = body_pr.find('a:normAutofit', nsmap)
            if norm is None:
                continue
            font_scale = int(norm.get('fontScale', '100000')) / 100000.0
            if font_scale < 1.0:
                for r in shape.text_frame._txBody.findall('.//a:r', nsmap):
                    rpr = r.find('a:rPr', nsmap)
                    if rpr is not None and rpr.get('sz'):
                        rpr.set('sz', str(int(int(rpr.get('sz')) * font_scale)))
            body_pr.remove(norm)
            changed = True
    if changed:
        prs.save(str(pptx_path))


def _check_readability(pptx_path):
    """Check for readability issues using PowerPoint's actual auto-fit results.

    Must run AFTER _refresh_autofit (so PowerPoint has calculated fontScale)
    and BEFORE _unlock_height_constraints (which removes fontScale data).
    """
    prs = Presentation(str(pptx_path))
    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    issues = []

    EMU_PER_PX = 6350
    SAFE_BOTTOM_PX = 950
    OVERFLOW_TOLERANCE = 80
    MIN_FONT_SZ = 1050  # 10.5pt in hundredths of a point

    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1

        for shape in slide.shapes:
            try:
                bottom_px = (shape.top + shape.height) / EMU_PER_PX
                if shape.has_text_frame and bottom_px > SAFE_BOTTOM_PX + OVERFLOW_TOLERANCE:
                    text_preview = shape.text_frame.text[:30].replace('\n', ' ')
                    issues.append(
                        f"Slide {slide_num}: Shape extends to y={int(bottom_px)}px "
                        f"(safe area: ~{SAFE_BOTTOM_PX}px) — '{text_preview}'"
                    )
            except Exception:
                pass

            if not shape.has_text_frame:
                continue

            body_pr = shape.text_frame._txBody.find('.//a:bodyPr', nsmap)
            if body_pr is None:
                continue

            norm = body_pr.find('a:normAutofit', nsmap)
            if norm is None:
                continue

            font_scale = int(norm.get('fontScale', '100000')) / 100000.0
            if font_scale >= 1.0:
                continue

            for r in shape.text_frame._txBody.findall('.//a:r', nsmap):
                rpr = r.find('a:rPr', nsmap)
                if rpr is None or not rpr.get('sz'):
                    continue

                original_sz = int(rpr.get('sz'))
                effective_sz = int(original_sz * font_scale)

                if effective_sz < MIN_FONT_SZ:
                    t_elem = r.find('a:t', nsmap)
                    text_preview = (t_elem.text or '')[:20] if t_elem is not None else ''
                    issues.append(
                        f"Slide {slide_num}: Text shrunk to {effective_sz / 100:.1f}pt "
                        f"(was {original_sz / 100:.1f}pt, scale {font_scale * 100:.0f}%) "
                        f"— '{text_preview}'"
                    )
                    break

    # --- Check 3: Textbox-textbox overlap (after PowerPoint has calculated actual sizes) ---
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    OVERLAP_MARGIN = 5  # Allow 5px overlap tolerance
    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        textboxes = []
        containers = []

        for shape in slide.shapes:
            try:
                x = shape.left / EMU_PER_PX
                y = shape.top / EMU_PER_PX
                w = shape.width / EMU_PER_PX
                h = shape.height / EMU_PER_PX
            except Exception:
                continue

            text = ''
            if shape.has_text_frame:
                text = shape.text_frame.text[:25].replace('\n', ' ')

            entry = {'x': x, 'y': y, 'w': w, 'h': h, 'text': text}

            # TEXT_BOX (17) = textbox added via add_textbox
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX and text.strip():
                textboxes.append(entry)
            # AUTO_SHAPE (1) = shape added via add_shape
            # Only count as container if it has a visible fill (not "none"/transparent)
            elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                try:
                    sp_pr = shape._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
                    grad_pr = shape._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}gradFill')
                    if sp_pr is not None or grad_pr is not None:
                        containers.append(entry)
                except Exception:
                    pass

        # Check textbox vs textbox overlap
        for i, tb1 in enumerate(textboxes):
            for j, tb2 in enumerate(textboxes):
                if i >= j:
                    continue
                if (tb1['x'] < tb2['x'] + tb2['w'] - OVERLAP_MARGIN and
                    tb1['x'] + tb1['w'] > tb2['x'] + OVERLAP_MARGIN and
                    tb1['y'] < tb2['y'] + tb2['h'] - OVERLAP_MARGIN and
                    tb1['y'] + tb1['h'] > tb2['y'] + OVERLAP_MARGIN):
                    issues.append(
                        f"Slide {slide_num}: Textbox overlap — "
                        f"'{tb1['text']}' and '{tb2['text']}'"
                    )

        # Check textbox exceeding container bounds
        CONTAIN_MARGIN = 15  # top-left must be well inside, not on boundary
        for tb in textboxes:
            for ct in containers:
                # Is the textbox clearly inside this container? (not just touching the edge)
                if (ct['x'] + CONTAIN_MARGIN < tb['x'] < ct['x'] + ct['w'] - CONTAIN_MARGIN and
                    ct['y'] + CONTAIN_MARGIN < tb['y'] < ct['y'] + ct['h'] - CONTAIN_MARGIN):
                    # Check if textbox extends beyond container
                    if tb['x'] + tb['w'] > ct['x'] + ct['w'] + OVERLAP_MARGIN:
                        issues.append(
                            f"Slide {slide_num}: Textbox exceeds container width — "
                            f"'{tb['text']}' extends {int(tb['x'] + tb['w'] - ct['x'] - ct['w'])}px beyond card"
                        )
                    if tb['y'] + tb['h'] > ct['y'] + ct['h'] + OVERLAP_MARGIN:
                        issues.append(
                            f"Slide {slide_num}: Textbox exceeds container height — "
                            f"'{tb['text']}' extends {int(tb['y'] + tb['h'] - ct['y'] - ct['h'])}px below card"
                        )
                    break  # Only check the innermost container

    if issues:
        print(f"\n⚠️  Readability check ({len(issues)} issues):", file=sys.stderr)
        for issue in issues:
            print(f"  {issue}", file=sys.stderr)
        print("  → Shorten text, increase container, or split slides.\n", file=sys.stderr)
    else:
        print("✅ Readability check: Pass", file=sys.stderr)

    return issues


def cmd_generate(args):
    """Generate PPTX from JSON."""
    # Template path
    if args.template:
        template = Path(args.template)
    else:
        template = Path(__file__).parent.parent / "template_2026.pptx"
    
    if not template.exists():
        print(f"Error: Template not found: {template}", file=sys.stderr)
        sys.exit(1)
    
    # Read input
    if args.input and args.input != "-":
        data = json.loads(Path(args.input).read_text())
    else:
        data = json.load(sys.stdin)
    
    # Get theme
    theme = data.get("theme", "light")
    if theme not in LAYOUT_MAP:
        print(f"Error: Unknown theme: {theme}. Available: light, dark", file=sys.stderr)
        sys.exit(1)
    
    # Validate icons before generation
    missing_icons = validate_icons_in_json(data, theme)
    if missing_icons:
        print(f"Error: Missing icons ({len(missing_icons)}):", file=sys.stderr)
        for icon in sorted(missing_icons):
            print(f"  - icons:{icon}", file=sys.stderr)
        print(f"", file=sys.stderr)
        print(f"Run the following command to download icons:", file=sys.stderr)
        print(f"  python3 scripts/download_icons.py", file=sys.stderr)
        sys.exit(1)
    
    # Pre-generation JSON validation: warn about textboxes without height
    json_warnings = []
    for i, slide_def in enumerate(data.get("slides", []), 1):
        elements = slide_def.get("elements", [])
        for el in elements:
            if el.get("type") == "textbox" and el.get("text") and "height" not in el:
                # Check if there are elements below this textbox
                el_y = el.get("y", 0)
                has_elements_below = any(
                    other.get("y", 0) > el_y and other is not el
                    for other in elements
                    if other.get("type") in ("textbox", "shape", "table", "image")
                )
                if has_elements_below:
                    text_preview = el["text"][:20].replace("\n", " ")
                    json_warnings.append(
                        f"Slide {i}: textbox at y={el_y} has no height — "
                        f"may overlap elements below — '{text_preview}'"
                    )
            if el.get("type") == "textbox" and el.get("paragraphs") and "height" not in el:
                el_y = el.get("y", 0)
                has_elements_below = any(
                    other.get("y", 0) > el_y and other is not el
                    for other in elements
                    if other.get("type") in ("textbox", "shape", "table", "image")
                )
                if has_elements_below:
                    json_warnings.append(
                        f"Slide {i}: textbox (paragraphs) at y={el_y} has no height — "
                        f"may overlap elements below"
                    )
    if json_warnings:
        print(f"\n⚠️  JSON validation ({len(json_warnings)} warnings):", file=sys.stderr)
        for w in json_warnings:
            print(f"  {w}", file=sys.stderr)
        print("  → Add height to textboxes (1-line: fontSize×3.5, multi-line: lines×fontSize×2.7)\n", file=sys.stderr)

    # Generate
    builder = PPTXBuilder(template, theme)
    slides = data.get("slides", [])
    
    # Build id -> slide map for override resolution
    id_map = {}
    for slide_def in slides:
        if "id" in slide_def:
            sid = slide_def["id"]
            if sid in id_map:
                print(f"Error: Duplicate slide id: {sid}", file=sys.stderr)
                sys.exit(1)
            id_map[sid] = slide_def
    
    # Resolve overrides and add slides
    for slide_def in slides:
        resolved = resolve_override(slide_def, id_map)
        builder.add_slide(resolved)
    builder.save(Path(args.output))
    
    print(f"Generated: {args.output} (theme: {theme})")
    for i, slide_def in enumerate(slides, 1):
        title = slide_def.get("title", "(no title)")
        print(f"page{i:02d} - {title}")
    
    # Post-process: refresh autofit via PowerPoint, then unlock height constraints
    if not args.no_autofit:
        pptx_path = Path(args.output).resolve()
        _refresh_autofit(pptx_path)
        _check_readability(pptx_path)
        _unlock_height_constraints(pptx_path)


def cmd_preview(args):
    """Export PPTX slides as PNG images."""
    import subprocess
    import glob
    import tempfile
    
    pptx_path = Path(args.input).resolve()
    
    if args.output:
        output_dir = Path(args.output).resolve()
        output_dir.mkdir(parents=True, exist_ok=True)
    elif _is_wsl():
        # On WSL, use a directory next to the PPTX so PowerPoint COM can access it via Windows paths
        output_dir = pptx_path.parent / "preview"
        output_dir.mkdir(parents=True, exist_ok=True)
    else:
        output_dir = Path(tempfile.gettempdir()) / "pptx-preview"
        output_dir.mkdir(parents=True, exist_ok=True)
    
    if not pptx_path.exists():
        print(f"Error: File not found: {pptx_path}", file=sys.stderr)
        sys.exit(1)
    
    pdf_path = output_dir / "slides.pdf"
    
    # PPTX -> PDF: platform-specific
    is_wsl = _is_wsl()
    
    if is_wsl:
        win_pptx = subprocess.run(["wslpath", "-w", str(pptx_path)], capture_output=True, text=True).stdout.strip()
        win_pdf = subprocess.run(["wslpath", "-w", str(pdf_path)], capture_output=True, text=True).stdout.strip()
        ps_cmd = (
            f"$app = New-Object -ComObject PowerPoint.Application; "
            f"$prs = $app.Presentations.Open('{win_pptx}', "
            f"[Microsoft.Office.Core.MsoTriState]::msoTrue, "
            f"[Microsoft.Office.Core.MsoTriState]::msoFalse, "
            f"[Microsoft.Office.Core.MsoTriState]::msoFalse); "
            f"$prs.SaveAs('{win_pdf}', 32); "
            f"$prs.Close(); $app.Quit()"
        )
        result = subprocess.run(
            ["/mnt/c/Windows/System32/WindowsPowerShell/v1.0/powershell.exe", "-Command", ps_cmd],
            capture_output=True, timeout=60)
        # PowerShell on Windows outputs cp932 (Shift_JIS); decode with errors='replace' for cross-platform safety
        result = subprocess.CompletedProcess(
            result.args, result.returncode,
            stdout=result.stdout.decode("cp932", errors="replace") if result.stdout else "",
            stderr=result.stderr.decode("cp932", errors="replace") if result.stderr else ""
        )
    else:
        script = f'''
        tell application "Microsoft PowerPoint"
            open POSIX file "{pptx_path}"
            delay 2
            set theDoc to active presentation
            set outPath to (POSIX file "{pdf_path}") as text
            save theDoc in outPath as save as PDF
            close theDoc saving no
        end tell
        '''
        result = subprocess.run(["osascript", "-e", script], capture_output=True, text=True)
    
    if result.returncode != 0 or not pdf_path.exists():
        print(f"Error: PDF export failed: {result.stderr}", file=sys.stderr)
        sys.exit(1)
    
    # Parse pages to export
    pages_to_export = None
    if args.pages:
        pages_to_export = set(int(p.strip()) for p in args.pages.split(","))
    
    # PDF -> PNG via pdftoppm (export all, filter later)
    cmd = ["pdftoppm", "-png", str(pdf_path), str(output_dir / "page")]
    
    result = subprocess.run(cmd, capture_output=True, text=True)
    if result.returncode != 0:
        print(f"Error: PNG conversion failed: {result.stderr}", file=sys.stderr)
        sys.exit(1)
    
    # Rename with slide titles
    prs = Presentation(str(pptx_path))
    titles = {}
    for i, slide in enumerate(prs.slides, 1):
        title = ""
        if slide.shapes.title:
            title = slide.shapes.title.text.strip().replace("\n", " ")[:30]
        title = re.sub(r'[\\/:*?"<>|]', '', title)
        titles[i] = title or "notitle"
    
    generated = []
    for png in sorted(glob.glob(str(output_dir / "page-*.png"))):
        basename = Path(png).name
        match = re.match(r'page-(\d+)\.png', basename)
        if match:
            num = int(match.group(1))
            if pages_to_export and num not in pages_to_export:
                Path(png).unlink()  # Delete unwanted pages
                continue
            new_name = f"page{num:02d}-{titles.get(num, 'notitle')}.png"
            new_path = output_dir / new_name
            Path(png).rename(new_path)
            generated.append(new_path)
    
    # Cleanup PDF
    pdf_path.unlink()
    
    # Add grid overlay unless disabled
    if not args.no_grid:
        from PIL import Image, ImageDraw, ImageFont
        color = (255, 0, 0, 128)
        try:
            font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 24)
        except:
            try:
                font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 24)
            except:
                font = ImageFont.load_default()
        
        for png_path in generated:
            img = Image.open(png_path).convert("RGBA")
            w, h = img.size
            overlay = Image.new("RGBA", (w, h), (0, 0, 0, 0))
            draw = ImageDraw.Draw(overlay)
            for pct in range(5, 100, 5):
                x, y = int(w * pct / 100), int(h * pct / 100)
                px_x, px_y = int(1920 * pct / 100), int(1080 * pct / 100)
                draw.line([(x, 0), (x, h)], fill=color, width=1)
                draw.line([(0, y), (w, y)], fill=color, width=1)
                if pct % 10 == 0:
                    draw.text((x + 4, 4), f"{px_x}px ({pct}%)", fill=color, font=font)
                    draw.text((4, y + 4), f"{px_y}px ({pct}%)", fill=color, font=font)
            Image.alpha_composite(img, overlay).convert("RGB").save(png_path)
    
    for path in generated:
        print(f"Generated: {path}")


# Service name aliases for icon search
ICON_ALIASES = {
    "sqs": "simple queue service",
    "sns": "simple notification service",
    "ses": "simple email service",
    "s3": "simple storage service",
    "ebs": "elastic block store",
    "efs": "elastic file system",
    "elb": "elastic load balancing",
    "alb": "application load balancer",
    "nlb": "network load balancer",
    "ecr": "elastic container registry",
    "ecs": "elastic container service",
    "eks": "elastic kubernetes service",
    "vpc": "virtual private cloud",
    "iam": "identity and access management",
    "kms": "key management service",
    "acm": "certificate manager",
    "waf": "web application firewall",
    "ssm": "systems manager",
    "cfn": "cloudformation",
    "cdk": "cloud development kit",
    "cli": "command line interface",
    "ddb": "dynamodb",
    "msk": "managed streaming for apache kafka",
    "mwaa": "managed workflows for apache airflow",
    "emr": "elastic mapreduce",
}

def cmd_icon_search(args):
    """Search icons by keywords using manifest.json."""
    # Try Raycast manifest first, then local
    manifest_path = ICON_DIR.parent / "manifest.json"
    if not manifest_path.exists():
        manifest_path = ICON_LOCAL_DIR / "manifest.json"
    if not manifest_path.exists():
        _icons_not_installed_error()
    
    with open(manifest_path) as f:
        icons = json.load(f)["icons"]
    
    queries = args.query.lower().split()
    limit = args.limit
    type_filter = args.type
    theme_filter = args.theme
    
    # Search each query separately
    for query in queries:
        matches = []
        query_nospace = query.replace(" ", "").replace("-", "").replace("_", "")
        # Expand alias if exists
        expanded = ICON_ALIASES.get(query_nospace, "").replace(" ", "")
        for icon in icons:
            name_lower = icon["name"].lower()
            name_nospace = name_lower.replace(" ", "").replace("-", "").replace("_", "")
            if query_nospace in name_nospace or (expanded and expanded in name_nospace):
                if type_filter and icon["type"] != type_filter:
                    continue
                if not type_filter and icon["type"] == "shape":
                    continue
                if theme_filter:
                    file_lower = icon["file"].lower()
                    has_light = "_light" in file_lower or " light" in name_lower
                    has_dark = "_dark" in file_lower or " dark" in name_lower
                    if theme_filter == "light" and has_dark and not has_light:
                        continue
                    if theme_filter == "dark" and has_light and not has_dark:
                        continue
                # service type first, then shorter names
                type_priority = 0 if icon["type"] == "service" else 1
                score = (type_priority, len(icon["name"]))
                matches.append((score, icon))
        matches.sort(key=lambda x: (x[0], x[1]["name"]))
        
        print(f"# {query}")
        for _, icon in matches[:limit]:
            ratio = icon.get("aspectRatio", 1)
            h = int(100 / ratio) if ratio > 0 else 100
            size_str = "" if h == 100 else f"  (w:100 → h:{h})"
            label = f"{icon['name']} [{icon['category']}/{icon['type']}]"
            print(f"  {label:<60}{size_str}")
            print(f"    icons:{icon['file'].rsplit('.', 1)[0]}")
        if len(matches) > limit:
            print(f"  ... and {len(matches) - limit} more")
        print()


def _get_frontmatter(path: Path) -> dict:
    """Extract YAML frontmatter fields."""
    content = path.read_text()
    result = {"description": "", "category": ""}
    match = re.search(r'^---\s*\n(.*?)\n---', content, re.DOTALL)
    if match:
        fm = match.group(1)
        for key in result:
            m = re.search(rf'^{key}:\s*(.+)$', fm, re.MULTILINE)
            if m:
                result[key] = m.group(1).strip().strip('"\'')
    return result


def cmd_examples(args):
    """List or show design pattern examples."""
    script_dir = Path(__file__).parent.parent
    examples_dir = script_dir / "references"
    
    if not examples_dir.exists():
        print(f"Examples directory not found: {examples_dir}", file=sys.stderr)
        return
    
    patterns = sorted(f.stem for f in examples_dir.rglob("*.md"))
    
    if args.names:
        for name in args.names:
            candidates = list(examples_dir.rglob(f"{name}.md"))
            if not candidates:
                print(f"# Pattern not found: {name}", file=sys.stderr)
                print(f"# Available: {', '.join(patterns)}", file=sys.stderr)
                continue
            pattern_file = candidates[0]
            print(f"# {pattern_file}")
            print(pattern_file.read_text())
            print()
    else:
        groups = {}
        for md in sorted(examples_dir.rglob("*.md")):
            fm = _get_frontmatter(md)
            cat = fm["category"]
            groups.setdefault(cat, []).append((md.stem, fm["description"]))
        print("# Design Examples")
        print(f"# Path: {examples_dir}")
        for cat in ["component", "pattern"]:
            if cat not in groups:
                continue
            print(f"\n## {cat.title()}s")
            for name, desc in groups[cat]:
                print(f"  {name:<30} {desc}")


def _get_documents_dir():
    """Get Documents directory. On WSL2, use Windows Documents folder."""
    import subprocess
    if _is_wsl():
        try:
            result = subprocess.run(
                ["/mnt/c/Windows/System32/WindowsPowerShell/v1.0/powershell.exe", "-Command",
                 "[Environment]::GetFolderPath('MyDocuments')"],
                capture_output=True, timeout=10)
            win_path = result.stdout.decode("cp932", errors="replace").strip()
            if win_path:
                wsl = subprocess.run(["wslpath", win_path], capture_output=True, text=True)
                if wsl.returncode == 0:
                    return Path(wsl.stdout.strip())
        except Exception:
            pass
    return Path.home() / "Documents"


def cmd_init(args):
    from datetime import datetime
    if args.output:
        out_dir = Path(args.output).expanduser()
    else:
        ts = datetime.now().strftime("%Y%m%d-%H%M")
        name = f"{ts}-{args.name}" if args.name else ts
        out_dir = _get_documents_dir() / "pptx-maker" / name
    out_dir.mkdir(parents=True, exist_ok=True)
    json_path = out_dir / "presentation.json"
    json_path.write_text(json.dumps({"theme": args.theme, "slides": []}, indent=2) + "\n")
    print(f"output_dir:  {out_dir}")
    print(f"output_json: {json_path}")
    print(f"output_pptx: {out_dir / 'presentation.pptx'}")
    print(f"theme:       {args.theme}")


def main():
    parser = argparse.ArgumentParser(description="PPTX Builder")
    subparsers = parser.add_subparsers(dest="command", required=True)
    
    # generate
    p_gen = subparsers.add_parser("generate", help="Generate PPTX from JSON")
    p_gen.add_argument("input", nargs="?", help="Input JSON file (or - for stdin)")
    p_gen.add_argument("-o", "--output", required=True, help="Output PPTX path")
    p_gen.add_argument("-t", "--template", help="Template PPTX path")
    p_gen.add_argument("--no-autofit", action="store_true", help="Skip PowerPoint autofit refresh")
    
    # preview
    p_prev = subparsers.add_parser("preview", help="Export slides as PNG images")
    p_prev.add_argument("input", help="Input PPTX file")
    p_prev.add_argument("-o", "--output", help="Output directory (default: temp)")
    p_prev.add_argument("-p", "--pages", help="Pages to export (e.g. 1,3,5)")
    p_prev.add_argument("--no-grid", action="store_true", help="Disable 5% grid overlay")
    
    # icon-search
    p_icon = subparsers.add_parser("icon-search", help="Search AWS icons")
    p_icon.add_argument("query", help="Search keywords (space-separated)")
    p_icon.add_argument("-n", "--limit", type=int, default=20, help="Max results (default: 20)")
    p_icon.add_argument("-t", "--type", choices=["service", "resource", "group", "category", "general", "third-party"], help="Filter by icon type")
    p_icon.add_argument("--theme", choices=["light", "dark"], help="Filter by theme (light/dark)")
    
    # examples
    p_ex = subparsers.add_parser("examples", help="List or show design pattern examples")
    p_ex.add_argument("names", nargs="*", help="Pattern names to show (multiple allowed)")
    
    # init
    p_init = subparsers.add_parser("init", help="Initialize output directory with empty presentation JSON")
    p_init.add_argument("name", nargs="?", help="Presentation name (e.g. 'my-proposal')")
    p_init.add_argument("-o", "--output", help="Output directory (overrides default)")
    p_init.add_argument("--theme", choices=["light", "dark"], default="dark", help="Theme (default: dark)")
    
    args = parser.parse_args()
    
    if args.command == "generate":
        cmd_generate(args)
    elif args.command == "preview":
        cmd_preview(args)
    elif args.command == "icon-search":
        cmd_icon_search(args)
    elif args.command == "examples":
        cmd_examples(args)
    elif args.command == "init":
        cmd_init(args)

if __name__ == "__main__":
    main()
